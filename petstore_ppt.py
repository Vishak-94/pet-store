#!/usr/bin/env python3
"""
Generate a PowerPoint deck documenting the Java Pet Store 1.3.1_02 existing
architecture (low-level design). Embeds the Graphviz class diagram.

Output: PetStore_Architecture_LLD.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = "/Users/vishakvj/Downloads/pet-project"
DIAGRAM = os.path.join(BASE, "petstore_lld.png")

# ── palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x33, 0x55)
BLUE   = RGBColor(0x2E, 0x6D, 0xB4)
AMBER  = RGBColor(0xC9, 0x7C, 0x2F)
GREEN  = RGBColor(0x2F, 0x8F, 0x46)
RED    = RGBColor(0xB2, 0x3A, 0x48)
PURPLE = RGBColor(0x6B, 0x4F, 0xA0)
GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def textbox(s, x, y, w, h, lines, size=18, color=NAVY, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, tuple):
            text, opts = ln
        else:
            text, opts = ln, {}
        r = p.add_run(); r.text = text
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.color.rgb = opts.get("color", color)
        r.font.name = opts.get("font", font)
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])
        p.level = opts.get("level", 0)
    return tb


def title_bar(s, title, subtitle=None):
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    textbox(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.7),
            title, size=28, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(s, Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.35),
                subtitle, size=13, color=RGBColor(0xC5, 0xD3, 0xE6))


def bullets(s, x, y, w, h, items, size=16, gap=6):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0; text = it; color = NAVY; bold = False; sz = size
        if isinstance(it, tuple):
            text, meta = it
            lvl = meta.get("level", 0); color = meta.get("color", NAVY)
            bold = meta.get("bold", False); sz = meta.get("size", size)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(gap)
        r = p.add_run()
        bullet = "•  " if lvl == 0 else "–  "
        r.text = bullet + text
        r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Calibri"
    return tb


def table(s, x, y, w, h, data, col_widths=None, header_color=NAVY, fsize=11):
    rows, cols = len(data), len(data[0])
    gt = s.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gt.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = gt.cell(r, c)
            cell.text = str(data[r][c])
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(fsize)
            para.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            cell.margin_left = Pt(5); cell.margin_right = Pt(5)
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_color
                para.font.color.rgb = WHITE; para.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
                para.font.color.rgb = NAVY
    return gt


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.7), SW, Inches(0.06), AMBER)
textbox(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.3),
        "Java Pet Store 1.3.1_02", size=44, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(2.85), Inches(11.7), Inches(0.8),
        "Existing Architecture — Low-Level Design", size=26,
        color=RGBColor(0xC5, 0xD3, 0xE6))
textbox(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.2), [
    ("Reverse-engineered from source as the baseline for a modernization migration",
     {"size": 16, "color": WHITE, "space_after": 14}),
    ("J2EE 1.3 BluePrints  ·  4 enterprise apps  ·  19 components  ·  ~36k LOC  ·  EJB 2.x CMP  ·  JMS",
     {"size": 14, "color": RGBColor(0x9D, 0xB4, 0xD4)}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive summary + metrics
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Executive Summary", "What the system is")
bullets(s, Inches(0.5), Inches(1.35), Inches(6.6), Inches(5.6), [
    ("A distributed J2EE 1.3 e-commerce system (Sun BluePrints, 2003).", {"bold": True, "size": 17}),
    "Online pet store: catalog browse, cart, checkout, async order fulfilment, supplier integration, admin.",
    ("It is NOT one application —", {"bold": True, "color": RED, "size": 17}),
    ("4 separately-deployed .ear apps that talk asynchronously over JMS.", {"level": 1}),
    ("They share a library of 19 reusable EJB components.", {"level": 1}),
    ("Each app owns its own DB view; shared data travels as XML over JMS.", {"level": 1, "color": RED}),
    ("Persistence is split two ways:", {"bold": True, "size": 17}),
    ("EJB 2.x CMP (container-generated) for most entities.", {"level": 1}),
    ("One hand-written DAO (catalog) with a clean relational schema.", {"level": 1}),
])
table(s, Inches(7.4), Inches(1.5), Inches(5.4), Inches(4.8), [
    ["Metric", "Value"],
    ["Deployable apps (.ear)", "4 — petstore, opc, supplier, admin"],
    ["Reusable components", "19"],
    ["Java source files", "~309"],
    ["Lines of Java", "~36,000"],
    ["JSP pages", "98"],
    ["EJB descriptors", "20"],
    ["Persistence", "EJB 2.x CMP + 1 DAO"],
    ["Database", "Cloudscape (→ Derby)"],
    ["Messaging", "JMS point-to-point"],
    ["Runtime", "Sun J2EE 1.3 RI, JDK 1.4"],
], col_widths=[Inches(2.3), Inches(3.1)], fsize=12)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Application / deployment view
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Application (Deployment) View", "4 apps coupled only through JMS queues")
mono = (
    "Browser ─HTTP─▶  petstore.ear  (storefront)\n"
    "                   WAF MainServlet → EJBControllerLocalEJB (session façade)\n"
    "                   → catalog / cart / customer / signon components\n"
    "                   checkout builds PurchaseOrder (XML)\n"
    "                            │  jms/PurchaseOrderQueue\n"
    "                            ▼\n"
    "                 opc.ear  (Order Processing Center)\n"
    "                   PurchaseOrderMDB → ProcessManager (workflow state)\n"
    "                   large order → jms/OrderApprovalQueue ◀── admin approves\n"
    "                   InvoiceMDB, Mail MDBs → jms/*MailQueue → mailer → email\n"
    "                            │  jms/supplier  (forward PO as XML)\n"
    "                            ▼\n"
    "                 supplier.ear  (warehouse) — SupplierOrderMDB, Inventory\n"
    "                            │  invoice (XML) ──▶ back to opc\n"
    "                 petstoreadmin.ear  — Swing client (Java Web Start) → opc\n"
)
box = textbox(s, Inches(0.5), Inches(1.35), Inches(9.0), Inches(5.7),
              [(mono, {"size": 12.5, "font": "Consolas", "color": NAVY})])
rect(s, Inches(9.8), Inches(1.5), Inches(3.15), Inches(5.3), LIGHT)
textbox(s, Inches(9.95), Inches(1.6), Inches(2.9), Inches(5.1), [
    ("KEY FACT", {"bold": True, "color": RED, "size": 14, "space_after": 8}),
    ("The only edges between the 4 apps are JMS queues:", {"size": 13, "space_after": 6}),
    ("PurchaseOrderQueue", {"size": 12, "color": BLUE, "level": 1}),
    ("OrderApprovalQueue", {"size": 12, "color": BLUE, "level": 1}),
    ("supplier", {"size": 12, "color": BLUE, "level": 1}),
    ("*MailQueue", {"size": 12, "color": BLUE, "level": 1, "space_after": 10}),
    ("This is database-per-service / distributed data — 2003's version of microservice data isolation.",
     {"size": 13, "color": NAVY}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The class diagram (embedded)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Class-Level Design (by Bounded Context)",
          "Service (blue) · Domain (purple) · Data access (amber) · Relational table (green) · CMP table (red) · JMS (dashed)")
if os.path.exists(DIAGRAM):
    from PIL import Image
    iw, ih = Image.open(DIAGRAM).size
    avail_w = Inches(12.6); avail_h = Inches(5.9)
    scale = min(avail_w / iw, avail_h / ih)
    w = int(iw * scale); h = int(ih * scale)
    x = int((SW - w) / 2); y = Inches(1.25) + int((avail_h - h) / 2)
    s.shapes.add_picture(DIAGRAM, x, y, width=w, height=h)
else:
    textbox(s, Inches(0.5), Inches(3), Inches(12), Inches(1),
            "petstore_lld.png not found — run petstore_lld_diagram.py", color=RED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Service / business logic layer
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Service & Business-Logic Layer", "Session beans, MDBs, and the WAF controllers")
table(s, Inches(0.4), Inches(1.35), Inches(12.5), Inches(5.5), [
    ["Class", "Stereotype", "Key methods / responsibility"],
    ["MainServlet", "Front controller (Servlet)", "init(); doGet/doPost → doProcess()"],
    ["EJBControllerLocalEJB", "Session façade (stateful SB)", "processEvent(Event): EventResponse + StateMachine"],
    ["CatalogEJB", "Session bean", "getCategory, getProducts, getItem, searchItems"],
    ["ShoppingCartLocalEJB", "Stateful session bean", "addItem, updateItemQuantity, deleteItem, getSubTotal"],
    ["SignOnEJB", "Session bean", "authenticate(user,pwd): boolean; createUser"],
    ["ProcessManagerEJB", "Session bean (workflow)", "createManager, updateStatus, getStatus, getOrdersByStatus"],
    ["PurchaseOrderMDB", "Message-driven bean", "onMessage → persist PO, start workflow"],
    ["OrderApprovalMDB / InvoiceMDB", "Message-driven beans", "approval / invoice transitions"],
    ["SupplierOrderMDB", "Message-driven bean", "onMessage → fulfil PO, ship, invoice"],
], col_widths=[Inches(3.0), Inches(3.0), Inches(6.5)], fsize=12)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Domain model + data access
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Domain Model & Data-Access Layer")
textbox(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
        [("Domain value objects (persistence-agnostic)", {"bold": True, "size": 16, "color": PURPLE})])
bullets(s, Inches(0.6), Inches(1.75), Inches(12), Inches(1.5), [
    ("Catalog: Category, Product, Item, Page", {"size": 14}),
    ("Cart: Cart, CartItem", {"size": 14}),
    ("Order: PurchaseOrder, LineItem  — marshalled to/from XML for JMS transport", {"size": 14, "color": RED}),
])
textbox(s, Inches(0.5), Inches(3.25), Inches(12), Inches(0.4),
        [("Data-access — two mechanisms", {"bold": True, "size": 16, "color": AMBER})])
table(s, Inches(0.6), Inches(3.7), Inches(12.2), Inches(3.0), [
    ["Mechanism", "Where", "Components"],
    ["Hand-written DAO", "CatalogDAO (port) + Cloudscape/Generic adapters; SQL in CatalogDAOSQL.xml", "catalog"],
    ["EJB 2.x CMP", "Abstract entity beans; container generates SQL from <cmp-field> in ejb-jar.xml",
     "customer, account, profile, signon, creditcard, contactinfo, address, lineitem, purchaseorder, supplierpo, processmanager, uidgen"],
], col_widths=[Inches(2.3), Inches(6.4), Inches(3.5)], fsize=11)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Data model: two schema styles
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Data Model — Two Schema Styles", "The single biggest migration surface")
# left: relational
rect(s, Inches(0.4), Inches(1.35), Inches(6.1), Inches(0.5), GREEN)
textbox(s, Inches(0.5), Inches(1.4), Inches(5.9), Inches(0.4),
        [("Catalog — hand-designed relational (clean)", {"bold": True, "color": WHITE, "size": 14})])
bullets(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.8), [
    ("Normalized; locale-split (base + _details per locale).", {"size": 13}),
    ("Real foreign keys, proper types (DECIMAL(10,2)).", {"size": 13}),
    ("category → product → item hierarchy", {"size": 13, "color": GREEN, "bold": True}),
    ("category / category_details", {"level": 1, "size": 12}),
    ("product / product_details", {"level": 1, "size": 12}),
    ("item / item_details (price, attr1..5)", {"level": 1, "size": 12}),
    ("Migrates easily — already port/adapter shaped.", {"size": 13, "bold": True, "color": GREEN}),
])
# right: CMP
rect(s, Inches(6.8), Inches(1.35), Inches(6.1), Inches(0.5), RED)
textbox(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(0.4),
        [("Everything else — CMP container-generated", {"bold": True, "color": WHITE, "size": 14})])
bullets(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(4.8), [
    ("Auto-generated from <cmp-field> (sun-j2ee-ri.xml).", {"size": 13}),
    ("Tell-tale: __PMPrimaryKey, __reverse_* FKs, all VARCHAR(255).", {"size": 13}),
    ("Customer: User/Customer/Account/Profile/ContactInfo/Address/CreditCard", {"level": 1, "size": 12}),
    ("Order: PurchaseOrder/LineItem(+join)/Manager", {"level": 1, "size": 12}),
    ("Supplier: SupplierOrder/Inventory", {"level": 1, "size": 12}),
    ("ContactInfo/Address/CreditCard/LineItem DUPLICATED across all 3 apps.",
     {"size": 13, "color": RED, "bold": True}),
    ("Needs reverse-engineering + redesign + data-migration script.", {"size": 13, "bold": True}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Patterns → modern equivalents
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Legacy Patterns → Modern Equivalents", "What each 2003 construct becomes")
table(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.5), [
    ["Pet Store (2003)", "Modern equivalent (Spring Boot / Java 21)"],
    ["WAF MainServlet front controller", "Spring MVC DispatcherServlet"],
    ["EJBControllerLocalEJB + StateMachine + EJBAction", "@Controller / @Service methods"],
    ["EJB 2.x Local Home/Interface/Bean trios", "plain @Service / @Component beans"],
    ["CMP entity beans", "JPA entities + Spring Data repositories"],
    ["Catalog DAO + CatalogDAOSQL.xml", "Spring Data repository / JdbcTemplate adapter"],
    ["ServiceLocator (JNDI caching)", "Spring dependency injection (eliminated)"],
    ["JMS + MDBs between 4 EARs", "Spring events, or Kafka/RabbitMQ, or in-process (monolith)"],
    ["xmldocuments (XML over JMS)", "JSON DTOs / records"],
    ["JSP + custom taglibs + WAF templates", "Thymeleaf / React"],
    ["Cloudscape", "H2 / PostgreSQL / MongoDB"],
], col_widths=[Inches(6.0), Inches(6.3)], fsize=12)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Migration approach
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Migration Implications & Approach", "How to break up the work, de-risked")
bullets(s, Inches(0.5), Inches(1.4), Inches(12.2), Inches(5.5), [
    ("Strategy: Strangler Fig — stand up the new system beside the old behind a façade; migrate one bounded context at a time; never big-bang.", {"bold": True, "size": 16, "color": BLUE}),
    ("Sequence by bounded context (leaf/low-risk first):", {"bold": True, "size": 15}),
    ("1. Catalog — clean schema, no inbound deps → proves the stack", {"level": 1, "size": 14, "color": GREEN}),
    ("2. Customer / SignOn — self-contained, security-sensitive", {"level": 1, "size": 14}),
    ("3. Cart — session state, depends on catalog", {"level": 1, "size": 14}),
    ("4. Order / PurchaseOrder — the rich aggregate", {"level": 1, "size": 14}),
    ("5. OPC + Supplier fulfilment (JMS backbone) — migrate last", {"level": 1, "size": 14, "color": RED}),
    ("Principles that carry it:", {"bold": True, "size": 15}),
    ("Ports & adapters at every infra boundary — DB/messaging swaps become additive", {"level": 1, "size": 14}),
    ("Framework-free domain model (POJOs/records)", {"level": 1, "size": 14}),
    ("Contract tests pin the external surface; parallel-run to compare old vs new", {"level": 1, "size": 14}),
    ("Automate the mechanical (javax→jakarta, boilerplate); humans review the semantic", {"level": 1, "size": 14}),
    ("Data migration is its own workstream: map schema, backfill, reconcile, reversible scripts", {"level": 1, "size": 14}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Section divider: functional walkthrough
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.5), SW, Inches(0.06), AMBER)
textbox(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.0),
        "Functional Walkthrough", size=40, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.2), [
    ("Every use-case from the entry point — with the edge cases where it fails.", {"size": 18, "color": RGBColor(0xC5,0xD3,0xE6)}),
    ("This is the behavioural contract the characterization tests must pin before migrating.", {"size": 14, "color": RGBColor(0x9D,0xB4,0xD4)}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Entry point & request flow
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Entry Point & Request Flow", "Everything enters through ONE servlet — MainServlet (*.do)")
flow = (
    "Browser  GET/POST *.do\n"
    "   │\n"
    "   ▼\n"
    "MainServlet.doProcess()          ← single front controller\n"
    "   │  1. look up URL in mappings.xml\n"
    "   │  2. if useFlowHandler → run FlowHandler\n"
    "   │  3. RequestProcessor turns *.do into a WAF Event\n"
    "   ▼\n"
    "EJBControllerLocalEJB.processEvent()   (stateful session bean)\n"
    "   │  StateMachine maps Event → EJBAction (command)\n"
    "   ▼\n"
    "XxxEJBAction.perform(event)      ← business logic, calls component EJBs\n"
    "   │  returns EventResponse\n"
    "   ▼\n"
    "ScreenFlowManager                ← next JSP screen\n"
    "   │  on exception → exception-mapping → error screen\n"
    "   ▼\n"
    "JSP screen rendered back to browser"
)
textbox(s, Inches(0.5), Inches(1.35), Inches(7.6), Inches(5.7),
        [(flow, {"size": 12, "font": "Consolas", "color": NAVY})])
rect(s, Inches(8.4), Inches(1.5), Inches(4.5), Inches(0.5), RED)
textbox(s, Inches(8.5), Inches(1.55), Inches(4.3), Inches(0.4),
        [("Exception routes (the edge-case map)", {"bold": True, "color": WHITE, "size": 13})])
table(s, Inches(8.4), Inches(2.15), Inches(4.5), Inches(1.8), [
    ["Exception", "→ Screen"],
    ["ShoppingCartEmptyOrder", "cart_empty_order_error"],
    ["DuplicateAccount", "duplicate_account"],
    ["GeneralFailure", "error"],
], col_widths=[Inches(2.5), Inches(2.0)], fsize=10)
textbox(s, Inches(8.4), Inches(4.3), Inches(4.5), Inches(2.6), [
    ("URL→screen and Event→EJBAction wiring lives entirely in mappings.xml (data, not code).",
     {"size": 12, "color": NAVY, "space_after": 8}),
    ("No per-page servlets — one controller, declarative flow.", {"size": 12, "color": GREY}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Use cases 1-3: catalog / signon / register
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Use-Cases 1–3: Browse · Sign-On · Register", "Entry → behaviour → failure edge cases")
table(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(5.6), [
    ["Use-case (entry)", "Happy path", "Edge cases where it FAILS / surprises"],
    ["Browse catalog\nGET /category.do?id=CATS", "CatalogEJB→DAO returns locale rows, paginated (Page)",
     "Unknown id → EMPTY list, NOT 404 · start=9999 → empty · start<0 → SQL error · DB down → error.screen"],
    ["Sign on\nPOST /signon.do", "authenticate(user,pwd)=true → session signed-on",
     "Wrong pwd & unknown user → SAME false (no distinction) · no lockout/counter · userName case-sensitive (Jane≠jane)"],
    ["Register account\nPOST /createcustomer.do (FlowHandler)", "Creates Customer→Account→Profile→ContactInfo→Address graph",
     "Duplicate username → DuplicateAccountException → duplicate_account.screen · partial create rolls back (CMP tx) · weak field validation"],
], col_widths=[Inches(3.2), Inches(4.0), Inches(5.4)], fsize=11)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Use case 4: cart (richest edge cases)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Use-Case 4: Shopping Cart", "Richest edge cases — source-confirmed behaviour that must be preserved")
table(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(5.6), [
    ["Operation", "Actual code behaviour", "Gotcha / failure edge case"],
    ["ADD_ITEM", "cart.put(itemID, 1)", "Re-adding an item RESETS qty to 1 — does NOT increment. 'Fixing' this breaks the contract."],
    ["UPDATE_ITEMS qty>0", "remove then put(qty)", "Sets absolute quantity ✓"],
    ["UPDATE_ITEMS qty≤0", "remove, does NOT re-add", "SILENTLY deletes the line item. Not an error."],
    ["UPDATE unknown itemId", "put a line for it", "Item not in catalog → getSubTotal later throws (dangling ref)"],
    ["getSubTotal()", "loops items; calls CatalogHelper.getItem PER line", "Item removed from catalog but still in cart → CatalogException → subtotal fails"],
    ["Session timeout", "stateful bean destroyed", "Cart silently empties. Migration to @SessionScope must reproduce; singleton would leak carts across users"],
], col_widths=[Inches(2.6), Inches(4.2), Inches(5.8)], fsize=10.5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Use cases 5-6: checkout + fulfilment (async)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Use-Cases 5–6: Checkout & Fulfilment", "The async JMS workflow — at-least-once delivery")
table(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(5.6), [
    ["Step (entry)", "Behaviour", "Edge case where it fails"],
    ["Checkout\nPOST /order.do", "PO built from cart+customer → XML → jms/PurchaseOrderQueue → order_complete.screen",
     "Empty cart → ShoppingCartEmptyOrderException → error screen · not signed on → blocked by filter"],
    ["OPC consumes PO", "PurchaseOrderMDB → ProcessManager status PENDING; large order → OrderApprovalQueue",
     "JMS broker down → enqueue fails, NO order created · duplicate redelivery → MDB must be IDEMPOTENT (no double order)"],
    ["Supplier fulfilment", "SupplierOrderMDB decrements InventoryEJB, ships, returns Invoice (XML)",
     "Concurrent orders for last unit → inventory race · Invoice for unknown orderId → no manager row → silent drop"],
    ["Notifications", "Mail MDBs send confirmation email", "Mail server down → redelivery → must not double-send"],
], col_widths=[Inches(2.6), Inches(4.6), Inches(5.4)], fsize=10)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Use cases 7-8 + cross-cutting failures
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Use-Cases 7–8 & Cross-Cutting Failures")
textbox(s, Inches(0.5), Inches(1.25), Inches(12), Inches(0.4),
        [("Locale · Sign-off · Admin", {"bold": True, "size": 15, "color": BLUE})])
table(s, Inches(0.4), Inches(1.7), Inches(12.5), Inches(2.1), [
    ["Use-case", "Behaviour", "Edge case"],
    ["changelocale.do?locale=ja_JP", "FlowHandler sets session locale; re-renders", "Unsupported locale (no _details rows) → blank/fallback text"],
    ["signoff.do", "Invalidate session → signoff.screen", "Cart LOST on signoff (stateful bean destroyed) — intended"],
    ["Admin (petstoreadmin.ear)", "Swing/JWS client → approve orders via OPC", "Java Web Start WON'T launch on modern JVMs · re-approve shipped order must be idempotent"],
], col_widths=[Inches(3.2), Inches(4.3), Inches(5.0)], fsize=10.5)
textbox(s, Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
        [("Cross-cutting failure modes (every use-case)", {"bold": True, "size": 15, "color": RED})])
bullets(s, Inches(0.6), Inches(4.45), Inches(12.2), Inches(2.6), [
    ("DB (Cloudscape) unreachable → *DAOSysException / EJBException → GeneralFailureException → error.screen", {"size": 12}),
    ("JMS broker down → async steps stall; user-facing enqueue → error.screen", {"size": 12}),
    ("Session loss/timeout → stateful cart & sign-on vanish → bounced to signon", {"size": 12}),
    ("Malformed XML message (xmldocuments) → MDB onMessage throws → JMS redelivery loop (poison message) unless a DLQ exists", {"size": 12}),
    ("Locale with no _details rows → blank/fallback text, never an error", {"size": 12}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — The dangerous edge cases (test-these-or-regress)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "The Dangerous Edge Cases", "The behaviours an agent (or human) will 'accidentally fix' — and break the contract")
bullets(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(4.2), [
    ("These MUST be pinned by characterization tests before migrating (CLAUDE.md §4, forbidden shortcuts #1 & #10):", {"bold": True, "size": 15, "color": NAVY}),
    ("addItem RESETS quantity to 1 — does not increment", {"level": 1, "size": 14, "color": RED, "bold": True}),
    ("updateItemQuantity(qty ≤ 0) silently DELETES the line", {"level": 1, "size": 14, "color": RED, "bold": True}),
    ("getSubTotal SILENTLY SKIPS items removed from the catalog (does NOT throw)", {"level": 1, "size": 14, "color": RED, "bold": True}),
    ("Unknown catalog id returns EMPTY, not 404", {"level": 1, "size": 14, "color": RED, "bold": True}),
    ("Checkout must be IDEMPOTENT under JMS redelivery (no double order)", {"level": 1, "size": 14, "color": RED, "bold": True}),
    ("Sign-on does NOT distinguish bad-password from unknown-user", {"level": 1, "size": 14, "color": RED, "bold": True}),
])
rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.95), LIGHT)
textbox(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8),
        [("Rule: pin behaviour AS-IS (including quirks). Migrate → run tests → green before commit. "
          "Never weaken a test to make it pass; never 'improve' behaviour during migration — parity first.",
          {"size": 13, "color": NAVY, "bold": True})], anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — COMPLETE API inventory (all 4 apps + web services + JMS)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Complete API Inventory (all 4 apps)",
          "Verified by full source scan — synchronous HTTP + SOAP endpoints")
table(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.7), [
    ["App", "Type", "Endpoint / API", "Purpose & edge case"],
    ["petstore", "HTTP *.do", "cart · order · signon · createcustomer · customer · changelocale · signoff",
     "Storefront (see use-case slides). Front controller MainServlet."],
    ["petstore", "Servlet", "/Populate", "Bootstraps catalog + accounts from XML; ~15 params (action, credit_card_*, forcefully). Re-run w/o forcefully → skips if data exists."],
    ["petstore", "Filter", "SignOnFilter (protected)", "Guards customer.do/.screen, enter_order_information.screen, signon_welcome.screen → redirect to signon, replay original URL."],
    ["opc", "SOAP (JAX-RPC)", "OPCService.submitInvoice(Source)", "Receives supplier invoice as XML. Bad XML → InvalidInvoiceException. This is a web-service variant of the app."],
    ["supplier", "SOAP (JAX-RPC)", "SupplierService.submitOrder(Source); queryOrderStatus(id)", "Receive PO / query status. Bad order → InvalidOrderException; unknown id → UnknownOrderIdException."],
    ["supplier", "Servlet", "/RcvrRequestProcessor · /Populate", "Inventory display & receive; seed inventory."],
    ["admin", "Servlet", "/AdminRequestProcessor · /ApplRequestProcessor", "Back-office order approval UI. Also a Swing/JWS client — WON'T launch on modern JVMs."],
], col_widths=[Inches(1.3), Inches(1.9), Inches(4.4), Inches(5.1)], fsize=10)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Async 'APIs': all 8 MDBs + JMS destinations
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Async APIs — 8 MDBs & JMS Destinations",
          "The message-driven surface (verified: 8 MDBs, queues + one Topic)")
table(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(4.3), [
    ["MDB (consumer)", "App", "Destination", "Responsibility"],
    ["PurchaseOrderMDB", "opc", "jms/PurchaseOrderQueue", "Consume checkout PO → ProcessManager PENDING; route large orders to approval"],
    ["OrderApprovalMDB", "opc", "jms/OrderApprovalQueue", "Apply admin approve/deny → next transition"],
    ["InvoiceMDB", "opc", "jms/InvoiceTopic (pub/sub)", "Supplier invoice → status SHIPPED. NOTE: a Topic (fan-out), not a queue"],
    ["MailCompletedOrderMDB", "opc", "jms/CompletedOrderMailQueue", "Email order-complete"],
    ["MailInvoiceMDB", "opc", "jms/MailQueue", "Email invoice"],
    ["MailOrderApprovalMDB", "opc", "jms/OrderApprovalMailQueue", "Email approval outcome"],
    ["SupplierOrderMDB", "supplier", "(ejb-jar dest)", "Fulfil PO, decrement InventoryEJB, ship, return invoice"],
    ["MailerMDB", "component", "jms/MailQueue", "Generic mail sender (shared)"],
], col_widths=[Inches(3.0), Inches(1.3), Inches(3.3), Inches(5.0)], fsize=9.5)
rect(s, Inches(0.35), Inches(5.75), Inches(12.6), Inches(1.15), LIGHT)
textbox(s, Inches(0.55), Inches(5.85), Inches(12.2), Inches(1.0),
        [("Async edge cases (all MDBs): at-least-once delivery → onMessage must be IDEMPOTENT (re-delivery must not double-process). "
          "Malformed XML → onMessage throws → redelivery loop / poison message unless a DLQ exists. "
          "InvoiceTopic fan-out means multiple subscribers may each react — order of side-effects matters.",
          {"size": 12, "color": NAVY, "bold": True})], anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Per-endpoint data flow: DB writes + JMS emissions
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Per-Endpoint Data Flow — DB & JMS",
          "What each endpoint WRITES and what it EMITS (source-verified)")
table(s, Inches(0.3), Inches(1.28), Inches(12.75), Inches(5.75), [
    ["Endpoint", "DB written", "JMS emitted", "Edge / note"],
    ["GET category/product/item/search.do", "— (reads catalog tables)", "none", "Unknown id → empty, not 404"],
    ["POST signon.do", "— (reads UserEJBTable)", "none", "authenticate() reads only"],
    ["POST createcustomer.do", "User, Customer, Account, Profile, ContactInfo, Address, CreditCard", "none", "Dup username → DuplicateAccount; whole graph rolls back (1 tx)"],
    ["POST customer.do (update)", "Account, ContactInfo, Address, CreditCard", "none", "Protected URL (sign-on required)"],
    ["cart.do (add/delete/update)", "NONE — in-memory session bean", "none", "No DB touch; lost on timeout/signoff"],
    ["POST order.do (checkout)", "NONE directly (Counter++ for order id)", "→ jms/AsyncSenderQueue (PO as XML)", "Empty cart → exception before send; order NOT persisted here; cart emptied after"],
    ["changelocale.do / signoff.do", "—", "none", "Session state only"],
    ["GET /Populate", "ALL catalog + sample account tables (bulk)", "none", "forcefully=true re-seeds; else skips"],
    ["SupplierService.submitOrder (SOAP)", "Inventory (--), SupplierOrder", "→ jms/opc/InvoiceTopic", "Bad order → InvalidOrderException"],
    ["OPCService.submitInvoice (SOAP)", "Manager (status)", "→ mail queue", "Bad XML → InvalidInvoiceException"],
    ["admin /AdminRequestProcessor", "reads Manager", "→ approval (drives OrderApprovalMDB)", "Re-approve shipped order must be idempotent"],
], col_widths=[Inches(3.4), Inches(3.5), Inches(2.9), Inches(2.95)], fsize=8.8)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — The async write-chain (where the order is actually written)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "The Async Write-Chain",
          "Order persistence is entirely async — no HTTP request writes order data")
chain = (
    "order.do ──JMS──▶ AsyncSenderQueue ──▶ jms/PurchaseOrderQueue\n"
    "  (writes: Counter++)                        │\n"
    "                                             ▼  PurchaseOrderMDB\n"
    "                    WRITES: PurchaseOrder, LineItem, ContactInfo,\n"
    "                            Address, CreditCard, Manager(status=PENDING)\n"
    "                                             │\n"
    "                     ┌───────────────────────┴─── large order?\n"
    "            yes ▼ jms/OrderApprovalQueue          no ▼ → supplier\n"
    "        OrderApprovalMDB                     SupplierOrderMDB\n"
    "        UPDATES: Manager(APPROVED)           WRITES: Inventory--, SupplierOrder\n"
    "                 │                                   │\n"
    "                 ▼ OrderApprovalMailQueue            ▼ jms/opc/InvoiceTopic  (pub/sub!)\n"
    "                                             InvoiceMDB\n"
    "                                             UPDATES: Manager(SHIPPED)\n"
    "                                                   │\n"
    "                                                   ▼ CompletedOrderMailQueue ─▶ MailQueue ─▶ email"
)
textbox(s, Inches(0.4), Inches(1.3), Inches(8.7), Inches(5.8),
        [(chain, {"size": 10.5, "font": "Consolas", "color": NAVY})])
rect(s, Inches(9.3), Inches(1.4), Inches(3.65), Inches(5.5), LIGHT)
textbox(s, Inches(9.45), Inches(1.5), Inches(3.4), Inches(5.3), [
    ("Migration-critical", {"bold": True, "color": RED, "size": 14, "space_after": 8}),
    ("Write-side is entirely async — checkout is fire-and-forget. Making it synchronous changes the consistency model (ADR).",
     {"size": 11.5, "space_after": 8}),
    ("ManagerEJBTable is written by 3 MDBs (create/approve/invoice) → consolidate behind one OrderStatusService.",
     {"size": 11.5, "space_after": 8}),
    ("Invoice is a TOPIC (fan-out), not a queue — switching to point-to-point changes delivery semantics.",
     {"size": 11.5, "space_after": 8}),
    ("Every MDB write is exposed to at-least-once redelivery → must be idempotent (characterization tests assert this).",
     {"size": 11.5}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Section divider: API split by actor
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.5), SW, Inches(0.06), AMBER)
textbox(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.0),
        "API Split-Up by Actor", size=40, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.2), [
    ("Every endpoint segregated by WHO uses it — the target service boundary.", {"size": 18, "color": RGBColor(0xC5,0xD3,0xE6)}),
    ("Customer · Supplier staff · Internal/Admin · System↔System · Ops", {"size": 14, "color": RGBColor(0x9D,0xB4,0xD4)}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Applications: who runs each
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "The Four Applications — Who Runs Each")
table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.2), [
    ["App (.ear)", "Primary user", "What it is"],
    ["petstore", "End customers (public)", "The storefront web app"],
    ["opc — Order Processing Center", "System (no human)", "Async order-workflow backend"],
    ["supplier", "Warehouse/supplier staff + OPC", "Fulfilment app + supplier integration"],
    ["admin", "Internal staff", "Back-office order approval console"],
], col_widths=[Inches(3.6), Inches(3.9), Inches(4.8)], fsize=13)
textbox(s, Inches(0.5), Inches(5.1), Inches(12.2), Inches(1.6), [
    ("Key: the actor segregation IS the target service boundary.", {"bold": True, "size": 15, "color": BLUE, "space_after": 6}),
    ("Only petstore is public. OPC is pure backend. Supplier & admin are staff/system-facing. "
     "Inter-app hops (JMS/SOAP) have no human user at all.", {"size": 13, "color": NAVY}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Customer-facing endpoints
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "1 · Customer-Facing Endpoints (petstore.ear)",
          "Public storefront — browse is anonymous; checkout/profile require login")
textbox(s, Inches(0.4), Inches(1.25), Inches(12), Inches(0.35),
        [("Browse (anonymous) — *.screen GET URLs via WAF + CatalogEJB", {"bold": True, "size": 13, "color": GREEN})])
table(s, Inches(0.4), Inches(1.62), Inches(12.5), Inches(1.9), [
    ["Endpoint", "Login?", "Role"],
    ["main.screen", "No", "Home page"],
    ["category.screen?category_id · product.screen?product_id · item.screen?item_id", "No", "Drill down catalog hierarchy"],
    ["search.screen?keywords", "No", "Search catalog"],
], col_widths=[Inches(7.0), Inches(1.3), Inches(4.2)], fsize=11)
textbox(s, Inches(0.4), Inches(3.75), Inches(12), Inches(0.35),
        [("Actions — *.do state changes via MainServlet", {"bold": True, "size": 13, "color": BLUE})])
table(s, Inches(0.4), Inches(4.12), Inches(12.5), Inches(2.7), [
    ["Endpoint", "Login?", "Role"],
    ["cart.do (add/delete/update)", "No", "Manage cart (in-memory, no DB)"],
    ["signon.do · createuser.do · createcustomer.do", "No", "Authenticate / register account"],
    ["changelocale.do · signoff.do", "No", "Switch language / log out"],
    ["customer.do", "YES", "View/edit own profile"],
    ["order.do (checkout)", "YES", "Place order → emits to JMS (no direct DB write)"],
], col_widths=[Inches(7.0), Inches(1.3), Inches(4.2)], fsize=11)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Supplier / Admin / Ops endpoints
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "2–3 · Supplier, Admin & Ops Endpoints")
textbox(s, Inches(0.4), Inches(1.25), Inches(12), Inches(0.35),
        [("Supplier-facing (supplier.ear)", {"bold": True, "size": 13, "color": AMBER})])
table(s, Inches(0.4), Inches(1.6), Inches(12.5), Inches(1.9), [
    ["Endpoint", "Actor", "Role"],
    ["/RcvrRequestProcessor", "Warehouse staff", "View pending supplier orders, trigger shipment"],
    ["SupplierService.submitOrder / queryOrderStatus (SOAP)", "OPC (system)", "Receive PO to fulfil / query status — web-services variant only"],
    ["/Populate (supplier)", "Ops", "Seed inventory"],
], col_widths=[Inches(5.4), Inches(2.6), Inches(4.5)], fsize=10.5)
textbox(s, Inches(0.4), Inches(3.7), Inches(12), Inches(0.35),
        [("Admin / back-office (admin.ear)", {"bold": True, "size": 13, "color": PURPLE})])
table(s, Inches(0.4), Inches(4.05), Inches(12.5), Inches(1.5), [
    ["Endpoint", "Actor", "Role"],
    ["/AdminRequestProcessor", "Internal staff", "Review & approve/deny large orders"],
    ["/ApplRequestProcessor + Swing/JWS client", "Internal staff", "Rich admin UI — WON'T launch on modern JVMs"],
], col_widths=[Inches(5.4), Inches(2.6), Inches(4.5)], fsize=10.5)
textbox(s, Inches(0.4), Inches(5.75), Inches(12), Inches(0.35),
        [("Ops / setup (one-time)", {"bold": True, "size": 13, "color": GREY})])
table(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.85), [
    ["Endpoint", "Actor", "Role"],
    ["/Populate (petstore)", "Operator", "Seed catalog + sample accounts from XML → becomes Flyway/data.sql"],
], col_widths=[Inches(5.4), Inches(2.6), Inches(4.5)], fsize=10.5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — System↔System integration seams
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "4 · System↔System Endpoints (no human user)",
          "The inter-app integration seams — customer never touches these")
table(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.0), [
    ["Channel / endpoint", "Producer → Consumer", "Transport", "Role"],
    ["AsyncSenderQueue → PurchaseOrderQueue", "petstore → OPC", "JMS", "Deliver checkout PO"],
    ["OrderApprovalQueue", "OPC → OPC", "JMS", "Route large order to approval"],
    ["jms/opc/InvoiceTopic", "Supplier → OPC", "JMS Topic", "Return invoice (order shipped) — fan-out"],
    ["jms/*MailQueue", "OPC → Mailer", "JMS", "Trigger notification emails"],
    ["OPCService.submitInvoice", "Supplier → OPC", "SOAP", "Return invoice — web-services variant of InvoiceTopic"],
    ["SupplierService.submitOrder", "OPC → Supplier", "SOAP", "Send PO — web-services variant of supplier queue"],
], col_widths=[Inches(4.2), Inches(3.0), Inches(1.8), Inches(3.5)], fsize=10.5)
rect(s, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.3), LIGHT)
textbox(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.1),
        [("The standard build uses JMS for these hops; the web-services build uses SOAP for the same flow. "
          "MIGRATE ONE TRANSPORT, NOT BOTH. These three contracts — submitOrder / submitInvoice / queryOrderStatus — "
          "become the inter-service API (REST or events).",
          {"size": 13, "color": NAVY, "bold": True})], anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Target service boundary (the split → new modules)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Target Split — Actor Boundary → New Modules",
          "The actor segregation becomes the Spring Boot module/service map")
table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.2), [
    ["Actor group", "Target module (Spring Boot / Java 21)", "Migration note"],
    ["Customer (public)", "Public REST/web module: browse + cart + checkout + account", "Migrate first after catalog; most traffic"],
    ["Supplier", "Fulfilment service with inbound 'receive order' endpoint", "submitOrder/queryOrderStatus = inter-service contract (REST or events)"],
    ["Admin", "Internal REST admin endpoints", "Swing/JWS client is REPLACED, not migrated"],
    ["System↔System", "In-process Spring events (monolith) OR message broker (services)", "Collapse JMS+SOAP hops; migrate ONE transport"],
    ["Ops/setup", "Flyway / data.sql seeding", "/Populate is NOT a runtime endpoint"],
], col_widths=[Inches(2.5), Inches(5.3), Inches(4.5)], fsize=11)
textbox(s, Inches(0.5), Inches(5.9), Inches(12.2), Inches(1.0),
        [("Bottom line: 'who uses it' defines the module boundary. Public/staff/system/ops separation "
          "maps cleanly onto bounded contexts — and tells you migration order (customer path first, admin last).",
          {"size": 13, "color": NAVY, "bold": True})])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 27 — How JMS communicates: the 3-part mechanism
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "How JMS Communicates Between Services",
          "Sender & receiver never know each other — coupled only by a logical destination name")
textbox(s, Inches(0.4), Inches(1.25), Inches(12), Inches(0.35),
        [("1 · Sender puts an XML message on a named destination (AsyncSenderEJB)", {"bold": True, "size": 13, "color": BLUE})])
textbox(s, Inches(0.5), Inches(1.62), Inches(12.2), Inches(1.55),
        [("qConnect = qFactory.createQueueConnection();\n"
          "session  = qConnect.createQueueSession(true, 0);   // true = transacted → at-least-once\n"
          "qSender  = session.createSender(q);\n"
          "TextMessage jmsMsg = session.createTextMessage();\n"
          "jmsMsg.setText(purchaseOrder.toXML());             // payload = PO as XML\n"
          "qSender.send(jmsMsg);                              // fire-and-forget, no reply",
          {"size": 11, "font": "Consolas", "color": NAVY})])
textbox(s, Inches(0.4), Inches(3.35), Inches(12), Inches(0.35),
        [("2 · Destination found via JNDI — never hardcoded (the decoupling)", {"bold": True, "size": 13, "color": GREEN})])
textbox(s, Inches(0.5), Inches(3.72), Inches(12.2), Inches(1.05),
        [('ASYNC_SENDER_QUEUE = "java:comp/env/jms/AsyncSenderQueue";   // logical alias\n'
          'q = serviceLocator.getQueue(ASYNC_SENDER_QUEUE);            // resolved at runtime\n'
          '// server config (setup.xml) maps alias → physical queue → binds the consuming MDB',
          {"size": 11, "font": "Consolas", "color": NAVY})])
textbox(s, Inches(0.4), Inches(4.95), Inches(12), Inches(0.35),
        [("3 · Receiver (MDB) reacts automatically — container calls onMessage()", {"bold": True, "size": 13, "color": AMBER})])
textbox(s, Inches(0.5), Inches(5.32), Inches(12.2), Inches(1.5),
        [("public void onMessage(Message recvMsg) {           // no polling — event-driven\n"
          "    String xml = ((TextMessage) recvMsg).getText();\n"
          "    approval = doWork(xml);       // parse PO, persist to DB, set status\n"
          "    if (approval != null) doTransition(approval);  // fire the NEXT message\n"
          "}   // binding declared in ejb-jar.xml: <resource-env-ref-name>jms/PurchaseOrderQueue",
          {"size": 11, "font": "Consolas", "color": NAVY})])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Two patterns + end-to-end conversation
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "JMS Patterns & End-to-End Conversation",
          "Both point-to-point (Queue) and publish/subscribe (Topic) are used")
table(s, Inches(0.4), Inches(1.35), Inches(6.0), Inches(1.9), [
    ["Pattern", "Class / call", "Used for"],
    ["Point-to-Point\n(Queue)", "QueueSender\n.send()", "PO, OrderApproval, Mail — ONE consumer"],
    ["Pub/Sub\n(Topic)", "TopicSender\n.publish()", "InvoiceTopic — EVERY subscriber (fan-out)"],
], col_widths=[Inches(1.7), Inches(1.9), Inches(2.4)], fsize=10)
conv = (
    "petstore              OPC                    supplier\n"
    "────────              ───                    ────────\n"
    "OrderAction\n"
    " send(PO xml)\n"
    "    │ AsyncSenderQueue\n"
    "    ▼\n"
    " [BROKER] ─▶ PurchaseOrderMDB.onMessage\n"
    "               persist PO, status=PENDING\n"
    "               send(PO) ──▶ supplier queue\n"
    "                              │\n"
    "                              ▼\n"
    "                     SupplierOrderMDB.onMessage\n"
    "                       Inventory--, ship\n"
    "                       publish(invoice)\n"
    "               InvoiceTopic ◀──┘ (pub/sub)\n"
    "  InvoiceMDB.onMessage\n"
    "    status=SHIPPED\n"
    "    ─▶ MailQueue ─▶ MailerMDB ─▶ email"
)
textbox(s, Inches(6.7), Inches(1.35), Inches(6.3), Inches(5.6),
        [(conv, {"size": 9.5, "font": "Consolas", "color": NAVY})])
textbox(s, Inches(0.4), Inches(3.6), Inches(6.0), Inches(3.3), [
    ("Why JMS (the properties it buys):", {"bold": True, "size": 13, "color": NAVY, "space_after": 6}),
    ("Decoupling — apps share only a queue name; redeploy independently", {"size": 12, "level": 1}),
    ("Async — checkout returns instantly; fulfilment happens later", {"size": 12, "level": 1}),
    ("Reliable — transacted session → redelivery if MDB throws (at-least-once)", {"size": 12, "level": 1}),
    ("Buffering — order spikes queue up; OPC drains at its own pace", {"size": 12, "level": 1}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 29 — JMS → Spring mapping (migration)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "JMS → Spring / Java 21 Mapping",
          "How the messaging translates in the migration")
table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.9), [
    ["Pet Store JMS", "Spring / Java 21 equivalent"],
    ["AsyncSenderEJB.send() (queue)", "JmsTemplate.convertAndSend() OR ApplicationEventPublisher.publishEvent()"],
    ["@MDB onMessage() (queue)", "@JmsListener OR @EventListener method"],
    ["TopicSender.publish() (topic)", "@JmsListener (topic factory) / Kafka topic / event to many listeners"],
    ["JNDI destination lookup + ServiceLocator", "application.yml connection config (ServiceLocator deleted)"],
    ["XML TextMessage payload", "JSON DTO / record"],
    ["transacted session (at-least-once)", "@Transactional listener + idempotency key"],
], col_widths=[Inches(5.3), Inches(7.0)], fsize=12)
rect(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.35), LIGHT)
textbox(s, Inches(0.7), Inches(5.65), Inches(11.9), Inches(1.15),
        [("ADR decision: keep it truly async with a broker (Kafka/RabbitMQ + @JmsListener) if preserving the distributed "
          "topology — OR collapse to in-process Spring ApplicationEvents for a modular monolith. Either way, preserve the "
          "Topic fan-out for invoices and at-least-once idempotency (both in the characterization-test spec).",
          {"size": 12.5, "color": NAVY, "bold": True})], anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 30 — High-level flow diagram (embedded)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "High-Level Flow — Apps, Endpoints & APIs",
          "Actors → apps → datastores · solid = sync (HTTP/SOAP) · dashed red = async (JMS)")
HL = os.path.join(BASE, "petstore_highlevel_flow.png")
if os.path.exists(HL):
    from PIL import Image
    iw, ih = Image.open(HL).size
    avail_w = Inches(12.8); avail_h = Inches(5.9)
    scale = min(avail_w / iw, avail_h / ih)
    w = int(iw * scale); h = int(ih * scale)
    x = int((SW - w) / 2); y = Inches(1.25) + int((avail_h - h) / 2)
    s.shapes.add_picture(HL, x, y, width=w, height=h)
else:
    textbox(s, Inches(0.5), Inches(3), Inches(12), Inches(1),
            "petstore_highlevel_flow.png not found — run petstore_highlevel_flow.py", color=RED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 31 — Migration plan: divider
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.5), SW, Inches(0.06), AMBER)
textbox(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.0),
        "Migration Plan", size=42, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.2), [
    ("Strangler-fig, bounded-context by bounded-context, contract-tested at every step.", {"size": 18, "color": RGBColor(0xC5,0xD3,0xE6)}),
    ("Target: Spring Boot 3.x · Java 21 · hexagonal (ports & adapters) · one transport, not both variants.", {"size": 14, "color": RGBColor(0x9D,0xB4,0xD4)}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 32 — Migration strategy & principles
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Migration Strategy & Principles")
bullets(s, Inches(0.5), Inches(1.4), Inches(12.2), Inches(5.6), [
    ("Strategy: Strangler Fig — new system beside old behind a façade; migrate one bounded context at a time; always shippable; easy rollback.", {"bold": True, "size": 15, "color": BLUE}),
    ("Characterization tests FIRST — pin legacy behaviour before touching code; migration becomes migrate → test → green-or-rollback.", {"bold": True, "size": 15, "color": GREEN}),
    ("Core principles:", {"bold": True, "size": 15}),
    ("Ports & adapters — DB & messaging behind interfaces → swaps become additive", {"level": 1, "size": 13}),
    ("Framework-free domain model (POJOs/records) — no JPA/Jackson on domain types", {"level": 1, "size": 13}),
    ("Preserve the contract — same status codes, message shapes, edge behaviours (see dangerous edge cases)", {"level": 1, "size": 13}),
    ("Migrate ONE transport, not both variants (JMS standard build, drop the SOAP duplicate)", {"level": 1, "size": 13}),
    ("Delete, don't port: ServiceLocator, Home/Local trios, WAF engine", {"level": 1, "size": 13, "color": RED}),
    ("Automate mechanical (javax→jakarta, boilerplate); humans review the semantic", {"level": 1, "size": 13}),
    ("Data migration is its own workstream: redesign CMP schema, backfill, reconcile, reversible scripts", {"level": 1, "size": 13}),
    ("State in a ledger (survives context resets) + ADRs for key decisions", {"level": 1, "size": 13}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 33 — Phased roadmap (by bounded context, actor-driven order)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Phased Roadmap — by Bounded Context",
          "Order = lowest-risk / most-value first; async backbone last")
table(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(5.6), [
    ["Phase", "Scope", "Why here / risk", "Done when"],
    ["0 · Scaffold", "Target repo, CI, CLAUDE.md, ledger, char-test harness", "De-risk before any rewrite", "Empty app deploys; harness runs"],
    ["1 · Catalog", "Browse APIs + catalog schema (clean DAO)", "No inbound deps; proves the stack", "Browse contract tests green"],
    ["2 · Customer/SignOn", "Register, profile, auth (CMP → JPA)", "Self-contained; security-sensitive", "Auth + profile tests green"],
    ["3 · Cart", "Cart APIs (stateful → @SessionScope)", "Depends on catalog; edge-case heavy", "Cart edge tests green"],
    ["4 · Order", "Checkout + PurchaseOrder aggregate", "Core domain; the rich write", "Checkout emits + persists correctly"],
    ["5 · OPC + Supplier", "MDB workflow, inventory, invoice (JMS backbone)", "Highest coupling; async + idempotency", "End-to-end order flow green"],
    ["6 · Admin", "Approval endpoints (REST replaces Swing/JWS)", "Swing client can't be carried forward", "Admin approval flow green"],
    ["7 · Cutover", "Data migration, parallel-run, decommission legacy", "Reconcile old↔new before switch", "Legacy off; ledger all-green"],
], col_widths=[Inches(2.2), Inches(3.9), Inches(3.4), Inches(3.1)], fsize=9.5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 34 — Target architecture & key decisions (ADRs)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Target Architecture & Key Decisions (ADRs)")
textbox(s, Inches(0.4), Inches(1.25), Inches(12.2), Inches(0.4),
        [("Target: modular monolith (default) — 1 Spring Boot app, module per bounded context, in-process events", {"bold": True, "size": 14, "color": BLUE})])
table(s, Inches(0.4), Inches(1.75), Inches(12.5), Inches(3.4), [
    ["Decision (ADR)", "Options", "Recommended default"],
    ["Async topology", "Keep broker (Kafka/Rabbit) vs in-process Spring events", "In-process events (monolith); broker only if OPC needs independent scale"],
    ["Transport variant", "JMS standard build vs SOAP web-services build", "Migrate JMS build; drop SOAP duplicate"],
    ["Persistence", "DB-per-context vs single consolidated schema", "Single schema (monolith); redesign CMP tables"],
    ["Checkout consistency", "Keep async fire-and-forget vs synchronous write", "Keep async (preserve semantics) — flag if UX needs sync"],
    ["Shared value objects", "Duplicate per context vs one shared domain lib", "Shared domain module (ContactInfo/Address/CreditCard/LineItem)"],
], col_widths=[Inches(3.0), Inches(4.8), Inches(4.7)], fsize=10)
rect(s, Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.5), LIGHT)
textbox(s, Inches(0.6), Inches(5.45), Inches(12.1), Inches(1.3),
        [("Migration order follows the high-level flow: customer path first (catalog→customer→cart→order), then the "
          "async OPC/Supplier backbone, admin last. Every phase: characterization tests green → commit → ledger update. "
          "Nothing is decommissioned until parallel-run reconciles old vs new.",
          {"size": 12.5, "color": NAVY, "bold": True})], anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 35 — Full database ER diagram (embedded)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Database Schema & Mappings — Full ER Diagram",
          "22 tables · green = relational (catalog) · red = CMP-generated · amber = join · ◆ PK ▸ FK")
ER = os.path.join(BASE, "petstore_er.png")
if os.path.exists(ER):
    from PIL import Image
    iw, ih = Image.open(ER).size
    avail_w = Inches(12.9); avail_h = Inches(6.0)
    scale = min(avail_w / iw, avail_h / ih)
    w = int(iw * scale); h = int(ih * scale)
    x = int((SW - w) / 2); y = Inches(1.2) + int((avail_h - h) / 2)
    s.shapes.add_picture(ER, x, y, width=w, height=h)
else:
    textbox(s, Inches(0.5), Inches(3), Inches(12), Inches(1),
            "petstore_er.png not found — run petstore_er_diagram.py", color=RED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 36 — Schema mapping notes (how the tables relate)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "How the Schema Maps — Key Relationships")
bullets(s, Inches(0.5), Inches(1.4), Inches(12.2), Inches(5.6), [
    ("Catalog (real FKs, locale-split):", {"bold": True, "size": 14, "color": GREEN}),
    ("category ─< product ─< item ; each base table has a _details child keyed by (id, locale) for i18n", {"level": 1, "size": 13}),
    ("Customer graph (CMP container-managed, via __PMPrimaryKey / __reverse_*):", {"bold": True, "size": 14, "color": RED}),
    ("User (login) ┈ Customer 1—1 Account 1—1 ContactInfo 1—1 Address ; Customer 1—1 Profile ; Account 1—1 CreditCard", {"level": 1, "size": 13}),
    ("Order graph:", {"bold": True, "size": 14, "color": RED}),
    ("PurchaseOrder 1—N LineItem (via join table) ; PO 1—1 ContactInfo (bill/ship) & CreditCard", {"level": 1, "size": 13}),
    ("ManagerEJBTable(orderId≈poId, status) = workflow spine — written by 3 MDBs", {"level": 1, "size": 13}),
    ("Supplier graph:", {"bold": True, "size": 14, "color": RED}),
    ("SupplierOrder 1—N LineItem (join) ; fulfils against Inventory via itemId", {"level": 1, "size": 13}),
    ("Duplicated shared tables — the big migration flag:", {"bold": True, "size": 14, "color": RED}),
    ("ContactInfo/Address (petstore+opc+supplier), CreditCard (petstore+opc), LineItem (opc+supplier) — each app has its own copy (database-per-app)", {"level": 1, "size": 13}),
    ("Target: drop __PMPrimaryKey/__reverse_*, fix types (poValue REAL→DECIMAL, poDate LONGINT→timestamp), consolidate shared tables into one domain module", {"level": 1, "size": 13, "color": BLUE, "bold": True}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 37 — Decision log: divider
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.5), SW, Inches(0.06), AMBER)
textbox(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.0),
        "Decision Log (ADRs)", size=42, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.2), [
    ("Every migration decision — the approaches we weighed and what we chose.", {"size": 18, "color": RGBColor(0xC5,0xD3,0xE6)}),
    ("DECIDED = locked by user · RECOMMENDED = proposed, awaiting confirm · OPEN = to be decided", {"size": 14, "color": RGBColor(0x9D,0xB4,0xD4)}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 38 — Decided
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Decisions — DECIDED (locked)",
          "Approaches considered → what we chose → why")
table(s, Inches(0.3), Inches(1.28), Inches(12.75), Inches(5.75), [
    ["Decision", "Approaches considered", "Chosen", "Why"],
    ["Source app", "kitchensink · Pet Store · both", "Pet Store only", "User focus; kitchensink out of scope"],
    ["Target runtime", "Spring Boot · Quarkus", "Spring Boot 3.x / Java 21", "Ecosystem, simple local run, exec JAR"],
    ["Async messaging", "Keep JMS (broker) · in-process events", "Keep JMS (@JmsListener, InvoiceTopic)", "User directive; faithful to distributed design"],
    ["SOAP variant", "Migrate it · drop it", "Drop — JMS build only", "Two builds are duplicates"],
    ["Database (goal 1)", "H2 · Postgres · MongoDB", "H2 embedded (Mongo = stretch)", "No install; runs on laptop"],
    ["Business logic", "Refactor/improve · preserve exactly", "Preserve exactly (char-tests)", "User directive; parity first"],
    ["Design approach", "Ad hoc · SOLID + patterns", "SOLID; ASK before choosing a pattern", "User directive"],
    ["Legacy code", "Edit in place · new project", "New petstore-app-v1/; legacy read-only", "It's a re-platform, not edits"],
], col_widths=[Inches(2.4), Inches(3.7), Inches(3.6), Inches(3.05)], fsize=9.5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 39 — Decisions resolved during implementation (phases 1-6)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Decisions — Resolved During Implementation",
          "Strategy/architecture + the pattern & technical calls made per phase")
table(s, Inches(0.3), Inches(1.28), Inches(12.75), Inches(5.75), [
    ["Decision", "Approaches considered", "Chosen"],
    ["Migration strategy", "Big-bang · Strangler Fig · lift-and-shift", "Strangler Fig, characterization-tests-first"],
    ["Architecture topology", "Modular monolith · microservices", "Modular monolith, ports & adapters"],
    ["Local layout", "Multi-Maven-module · single module + packages", "Single module, strict package boundaries"],
    ["JMS broker (local)", "Embedded Artemis · external ActiveMQ (Docker)", "Embedded Artemis"],
    ["Order workflow pattern", "GoF State pattern · enum + guarded transitions", "Enum + guarded transitions (right-sized)"],
    ["Web action dispatch", "Explicit Command pattern · thin controllers→services", "Thin controllers → services (idiomatic Spring)"],
    ["Inventory concurrency", "Atomic conditional UPDATE · @Version · pessimistic lock", "Pessimistic lock (SELECT…FOR UPDATE) — user-directed"],
    ["Checkout consistency", "Keep optimistic · real-time reservation", "Keep optimistic (async fulfilment)"],
    ["Admin client", "Migrate Swing/JWS · REST endpoints", "REST endpoints (Swing dead on modern JVM)"],
    ["Approve→fulfil trigger", "Publish inline in tx · AFTER_COMMIT event", "AFTER_COMMIT transactional event (race fix)"],
    ["Password storage", "Hash now · preserve plaintext (parity)", "Preserve plaintext (hashing = post-cutover)"],
], col_widths=[Inches(2.7), Inches(5.2), Inches(4.85)], fsize=9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 40 — Implementation status (what was actually built)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Implementation Status — petstore-app-v1",
          "Spring Boot 3.3 · Java 21 · embedded H2 + Artemis · 43 tests green")
table(s, Inches(0.4), Inches(1.35), Inches(12.5), Inches(3.9), [
    ["Phase", "Delivered", "Verified"],
    ["0 Scaffold", "Boot 3.3/Java 21, H2, embedded Artemis, char-test harness", "boots on Java 21"],
    ["1 Catalog", "domain, CatalogRepository port + JPA adapter, browse UI", "9 tests · browsable"],
    ["2 Customer/SignOn", "auth, registration, profile (CMP→JPA)", "7 tests"],
    ["3 Cart", "@SessionScope cart, all edge-case quirks preserved", "10 tests"],
    ["4 Order", "checkout, PO aggregate, enum workflow, JMS publish", "7 tests"],
    ["5 Fulfilment+JMS", "@JmsListener, pessimistic inventory, idempotent", "5 tests · live async"],
    ["6 Web+Admin", "REST admin approve/deny, exception handler", "5 tests · live e2e"],
    ["7 Cutover", "documented plan (data migration, parallel-run, decommission)", "PHASE7_CUTOVER.md"],
], col_widths=[Inches(2.6), Inches(7.1), Inches(2.8)], fsize=9.5)
rect(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.4), LIGHT)
textbox(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.2), [
    ("Two bugs caught by LIVE end-to-end testing (not unit tests): (1) large orders auto-shipped without approval; "
     "(2) approve→fulfil race (published before commit). Both fixed + regression-tested.", {"size": 12, "color": NAVY, "bold": True, "space_after": 4}),
    ("Ports keep DB (H2→Mongo) and broker (Artemis→Kafka) swaps additive — zero change to domain/services/tests.", {"size": 12, "color": BLUE}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 41 — Legacy add-to-cart & checkout flow
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Legacy: How a Customer Adds to Cart & Checks Out",
          "Reverse-engineered from the JSPs + mappings.xml + signon-config.xml")
flow = (
    "ADD TO CART (from item.jsp / product.jsp / search.jsp)\n"
    "  Each item row shows an \"Add to Cart\" link:\n"
    "    GET /cart.do?action=purchase&itemId=EST-1\n"
    "      → CartEvent(ADD_ITEM) → ShoppingCartLocalEJB.addItem()\n"
    "\n"
    "CHECKOUT (multi-step, sign-in gated)\n"
    "  1. cart.jsp → \"Proceed to Checkout\"\n"
    "        → enter_order_information.screen\n"
    "  2. That screen is a PROTECTED resource →\n"
    "        SignOnFilter forces sign-in first\n"
    "  3. enter_order_information.jsp = a FORM collecting\n"
    "        Billing info + Shipping info + Credit card\n"
    "        → POST order.do\n"
    "  4. order.do (OrderEJBAction) → build PO → JMS\n"
    "        → order_complete.screen"
)
textbox(s, Inches(0.5), Inches(1.35), Inches(7.7), Inches(5.6),
        [(flow, {"size": 12, "font": "Consolas", "color": NAVY})])
rect(s, Inches(8.5), Inches(1.5), Inches(4.4), Inches(5.0), LIGHT)
textbox(s, Inches(8.65), Inches(1.6), Inches(4.15), Inches(4.8), [
    ("Key legacy facts", {"bold": True, "color": RED, "size": 14, "space_after": 8}),
    ("Add-to-Cart link is on EVERY item — on the item page, product page, AND search results.", {"size": 12, "space_after": 8}),
    ("Checkout REQUIRES sign-in (SignOnFilter protects enter_order_information).", {"size": 12, "space_after": 8}),
    ("Checkout collects billing + shipping address + credit card on a form before placing the order.", {"size": 12, "space_after": 8}),
    ("The cart is a stateful session EJB — lost on sign-off / timeout.", {"size": 12}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 42 — Feature comparison: Legacy vs Migrated
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Feature Comparison — Legacy vs Migrated",
          "What's covered, and the remaining gaps")
table(s, Inches(0.3), Inches(1.28), Inches(12.75), Inches(5.75), [
    ["Feature", "Legacy", "Migrated", "Gap?"],
    ["Browse catalog (category→product→item)", "yes", "yes", "—"],
    ["Search", "yes", "yes", "—"],
    ["Add to Cart from item/product/search UI", "yes (link per item)", "backend only — NO UI link", "UI GAP"],
    ["Cart view / update / remove", "yes", "yes", "—"],
    ["Checkout requires sign-in", "yes", "yes (/checkout authenticated)", "—"],
    ["Checkout collects billing+shipping+card", "yes (form)", "only userId+email — no address/card capture", "GAP"],
    ["Order → async JMS fulfilment", "yes", "yes", "—"],
    ["Register / profile / address / credit card", "yes", "yes (REST endpoints)", "—"],
    ["Sign-on / logout", "yes (SignOnFilter)", "yes (Spring Security)", "—"],
    ["Admin approve / deny orders", "yes (Swing/JWS)", "yes (REST)", "—"],
    ["Inventory oversell → backorder", "yes", "yes (pessimistic lock)", "—"],
    ["i18n (en / ja / zh)", "yes (full)", "?lang= param; only en_US seeded", "Partial"],
    ["Admin sales charts (getChartInfo)", "yes", "not migrated", "Out of scope"],
], col_widths=[Inches(4.6), Inches(2.7), Inches(4.0), Inches(1.45)], fsize=8.5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 43 — Remaining gaps summary
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Remaining Gaps After Comparison")
textbox(s, Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
        [("Genuine functional gaps (worth closing for full parity):", {"bold": True, "size": 15, "color": RED})])
bullets(s, Inches(0.6), Inches(1.85), Inches(12.2), Inches(2.4), [
    ("Add-to-Cart UI link — legacy had it on every item/product/search row; migrated has the /cart/add endpoint but no button, so the store isn't click-to-buy in the browser.", {"size": 13}),
    ("Checkout form — legacy collects billing + shipping + credit card; migrated /checkout only takes userId+email and doesn't pull the saved account/card. Works, but thinner.", {"size": 13}),
    ("i18n — legacy fully multilingual (en/ja/zh); migrated supports ?lang= but only en_US data is seeded.", {"size": 13}),
])
textbox(s, Inches(0.5), Inches(4.4), Inches(12), Inches(0.4),
        [("Deliberately out of scope (documented):", {"bold": True, "size": 15, "color": GREY})])
bullets(s, Inches(0.6), Inches(4.9), Inches(12.2), Inches(1.6), [
    ("Admin sales charts (getChartInfo) — analytics", {"size": 13}),
    ("Supplier queryOrderStatus + warehouse inventory-view UI — internal supplier tooling", {"size": 13}),
    ("SOAP web-services variant — deliberately dropped (migrate JMS build only)", {"size": 13}),
])
rect(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.6), LIGHT)
textbox(s, Inches(0.7), Inches(6.42), Inches(11.9), Inches(0.5),
        [("Everything else — browse, cart, checkout, fulfilment, registration, account, card, sign-on/logout, admin — is migrated & tested (55 tests).",
          {"size": 12, "color": NAVY, "bold": True})], anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 44 — Section divider: customer-service
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.5), SW, Inches(0.06), AMBER)
textbox(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.0),
        "customer-service", size=44, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.2), [
    ("The extracted identity/accounts microservice — customer & auth endpoints only.", {"size": 18, "color": RGBColor(0xC5,0xD3,0xE6)}),
    ("Endpoints: POST /register · POST /auth/login · GET /customer/{id} · PUT /customer/{id}/account|profile|card",
     {"size": 13, "color": RGBColor(0x9D,0xB4,0xD4)}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 45 — Customer endpoints: legacy problems vs what we fixed
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Customer Endpoints — Legacy Problems vs Fixed",
          "Scope: register / login / customer profile-account-card only")
table(s, Inches(0.3), Inches(1.28), Inches(12.75), Inches(5.75), [
    ["Legacy problem (customer endpoints)", "Fixed?", "How"],
    ["Plaintext passwords (matchPassword = string equality)", "FIXED", "BCrypt hash at rest + encoder.matches"],
    ["Username = customer identity (guessable, enumerable)", "FIXED", "Opaque RANDOM customerId (stored UUID), in JWT cid claim"],
    ["Stateful sign-on (SignOnFilter + server session)", "FIXED", "Stateless JWT — issued on /auth/login, verified per request"],
    ["Customer PII + card duplicated / __PMPrimaryKey / VARCHAR(255)", "FIXED", "One owned customer schema, consolidated + typed columns"],
    ["Full card number returned on GET /customer", "FIXED", "cardMasked (**** **** **** 1111)"],
    ["No validation on register/update (blank/garbage accepted)", "FIXED", "@Valid — NotBlank userName, pwd 4-25, @Email; 400 field map"],
    ["Inconsistent error responses on customer endpoints", "FIXED", "@RestControllerAdvice — uniform {status,error,detail,correlationId}"],
    ["Hard to debug a login/register request", "FIXED", "Correlation-id per request (MDC + X-Correlation-Id header) + actuator"],
    ["No transport security on credentials (HTTP)", "DEFERRED", "TLS is gateway/infra — not fixed locally"],
], col_widths=[Inches(6.4), Inches(1.4), Inches(4.95)], fsize=9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 46 — Customer endpoints: further enhancements
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Customer Endpoints — Further Enhancements",
          "Next steps specific to auth + customer APIs")
textbox(s, Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
        [("Auth / credential security", {"bold": True, "size": 15, "color": RED})])
bullets(s, Inches(0.6), Inches(1.8), Inches(12.2), Inches(2.0), [
    ("TLS on /register and /auth/login so credentials aren't sent in clear", {"size": 13}),
    ("JWT: RS256 asymmetric signing + JWKS (so services verify without holding the mint key)", {"size": 13}),
    ("Refresh tokens + revocation denylist so /logout truly invalidates a token", {"size": 13}),
    ("Rate limiting + account lockout on /auth/login (brute-force defence)", {"size": 13}),
])
textbox(s, Inches(0.5), Inches(4.05), Inches(12), Inches(0.4),
        [("Customer data / ids", {"bold": True, "size": 15, "color": GREEN})])
bullets(s, Inches(0.6), Inches(4.5), Inches(12.2), Inches(2.0), [
    ("Snowflake-style customerId option: 64-bit, collision-free across DB shards without coordination + time-sortable (vs random UUID). Trade-off: leaks signup time/node — random UUID preferred for a customer id.", {"size": 13, "color": NAVY, "bold": True}),
    ("PCI: tokenize / vault the credit card instead of storing the PAN (even masked-on-read)", {"size": 13}),
    ("Email verification + password-reset flow on registration", {"size": 13}),
    ("Postgres for the customer DB + Flyway migrations (H2 today loses data on restart)", {"size": 13}),
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 47 — Decision rule: standalone service vs in-process library
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Service vs Library — When to Extract a Component",
          "The rule we applied when splitting the monolith (cart is a library, not a service)")

textbox(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
        [("Deciding question: does the component own remote, shared, or durable state — "
          "or need to scale / deploy independently? If yes → service. If no → in-process library.",
          {"bold": True, "size": 14, "color": NAVY})])

table(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.4), [
    ["Component", "Owns DB / shared durable state?", "Independent scale / deploy?", "Verdict"],
    ["customer-service", "Yes — accounts, credentials", "Yes (identity/JWT issuer)", "SERVICE"],
    ["catalog-service", "Yes — products (read-heavy)", "Yes (cacheable, read-scaled)", "SERVICE"],
    ["admin-office-service", "Yes — orders, inventory", "Yes (different traffic/roles)", "SERVICE"],
    ["cart (cart-lib)", "No — ephemeral per-session, in-memory", "No independent reason; on hot path", "LIBRARY"],
], col_widths=[Inches(2.6), Inches(4.2), Inches(3.6), Inches(1.9)], fsize=11)

textbox(s, Inches(0.5), Inches(4.55), Inches(6.1), Inches(0.4),
        [("Why cart is a LIBRARY (in-process)", {"bold": True, "size": 14, "color": GREEN})])
bullets(s, Inches(0.6), Inches(5.0), Inches(6.1), Inches(2.2), [
    ("Session-local, ephemeral state — no DB, no reason to be remote", {"size": 12}),
    ("A standalone server = a network hop for a user's own scratchpad (pure latency tax)", {"size": 12}),
    ("Stateful in-memory service is single-instance (or forces Redis) — extra distributed-state cost for no gain", {"size": 12}),
    ("On the checkout hot path; not an independently-evolving bounded context", {"size": 12}),
])

textbox(s, Inches(6.9), Inches(4.55), Inches(5.9), Inches(0.4),
        [("When cart would FLIP to a service", {"bold": True, "size": 14, "color": AMBER})])
bullets(s, Inches(7.0), Inches(5.0), Inches(5.8), Inches(2.2), [
    ("Carts must survive across devices / sessions (durable)", {"size": 12}),
    ("Storefront scales horizontally with carts shared across instances", {"size": 12}),
    ("→ then back it with Redis / a DB and extract — because it now owns SHARED state, not \"for consistency\"", {"size": 12, "color": NAVY, "bold": True}),
    ("\"Client\" was the wrong frame: a client is a thin proxy to a remote owner; cart IS the state + logic", {"size": 12}),
])

prs.save(os.path.join(BASE, "PetStore_Architecture_LLD.pptx"))
print("Wrote", os.path.join(BASE, "PetStore_Architecture_LLD.pptx"))
print("Slides:", len(prs.slides._sldIdLst))
