#!/usr/bin/env python3
"""
Java Pet Store — 45-MINUTE DEMO DECK (with code walkthrough).

A focused, presentation-paced deck for a 45-minute playback session:
  · the brief + what "done" means
  · gaps in the legacy design → what we fixed (parity audit)
  · a multi-step migration plan (incl. the legacy JMS XML → JSON plan)
  · target as-built architecture
  · a live-code walkthrough (real files, real snippets)
  · MongoDB stretch

It EMBEDS the already-generated diagrams (run these first if missing):
  petstore_legacy_highlevel.py   → petstore_legacy_highlevel.png
  petstore_architecture.py       → petstore_architecture.png
  petstore_packages.py           → petstore_packages.png

Output: docs/PetStore_Demo_45min.pptx   (self-contained; own copy of the helpers)
Deps: python-pptx (+ the PNGs above for the image slides).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

BASE = os.path.dirname(os.path.abspath(__file__))

# ── palette (shared with petstore_ppt.py) ─────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x33, 0x55)
BLUE   = RGBColor(0x2E, 0x6D, 0xB4)
AMBER  = RGBColor(0xC9, 0x7C, 0x2F)
GREEN  = RGBColor(0x2F, 0x8F, 0x46)
RED    = RGBColor(0xB2, 0x3A, 0x48)
PURPLE = RGBColor(0x6B, 0x4F, 0xA0)
GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF7)
CODEBG = RGBColor(0x1E, 0x28, 0x3B)   # dark panel for code
CODEFG = RGBColor(0xE6, 0xED, 0xF3)
CODEKW = RGBColor(0x7F, 0xB5, 0xF5)   # keyword-ish accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SUB    = RGBColor(0xC5, 0xD3, 0xE6)

# per-service accent colours (match the architecture diagram borders)
STORE_ACCENT = BLUE
AUTH_ACCENT  = RGBColor(0x6B, 0x5B, 0x3E)   # olive
CUST_ACCENT  = PURPLE
CAT_ACCENT   = RGBColor(0x2E, 0x86, 0xAB)   # teal-blue
OPC_ACCENT   = RGBColor(0xB8, 0x86, 0x0B)   # dark gold
INV_ACCENT   = AMBER
ADMIN_ACCENT = RED
NOTE_ACCENT  = GREEN

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def textbox(s, x, y, w, h, lines, size=18, color=NAVY, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, opts = ln if isinstance(ln, tuple) else (ln, {})
        r = p.add_run(); r.text = text
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.color.rgb = opts.get("color", color)
        r.font.name = opts.get("font", font)
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])
        p.level = opts.get("level", 0)
    return tb


def title_bar(s, title, subtitle=None, tag=None):
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    if tag:  # small amber "minute budget" chip on the right
        rect(s, Inches(11.35), Inches(0.30), Inches(1.75), Inches(0.55), AMBER)
        textbox(s, Inches(11.35), Inches(0.30), Inches(1.75), Inches(0.55),
                tag, size=13, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(0.5), Inches(0.12), Inches(10.7), Inches(0.7),
            title, size=27, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(s, Inches(0.5), Inches(0.74), Inches(10.7), Inches(0.35),
                subtitle, size=13, color=SUB)


def bullets(s, x, y, w, h, items, size=16, gap=7):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0; text = it; color = NAVY; bold = False; sz = size
        if isinstance(it, tuple):
            text, meta = it
            lvl = meta.get("level", 0); color = meta.get("color", NAVY)
            bold = meta.get("bold", False); sz = meta.get("size", size)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(gap)
        r = p.add_run()
        r.text = ("•  " if lvl == 0 else "–  ") + text
        r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Calibri"
    return tb


def table(s, x, y, w, h, data, col_widths=None, header_color=NAVY, fsize=11):
    rows, cols = len(data), len(data[0])
    gt = s.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gt.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = gt.cell(r, c)
            cell.text = str(data[r][c])
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(fsize); para.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            cell.margin_left = Pt(5); cell.margin_right = Pt(5)
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_color
                para.font.color.rgb = WHITE; para.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
                para.font.color.rgb = NAVY
    return gt


def code_panel(s, x, y, w, h, path, lines, caption=None):
    """A dark code panel: file-path header + monospaced body."""
    rect(s, x, y, w, h, CODEBG)
    # file path header
    hdr = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.06), w - Inches(0.3), Inches(0.3))
    hp = hdr.text_frame.paragraphs[0]
    hr = hp.add_run(); hr.text = path
    hr.font.size = Pt(11); hr.font.bold = True; hr.font.name = "Consolas"
    hr.font.color.rgb = CODEKW
    # body
    body = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.42), w - Inches(0.3), h - Inches(0.5))
    tf = body.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        text, opts = ln if isinstance(ln, tuple) else (ln, {})
        r = p.add_run(); r.text = text if text else " "
        r.font.size = Pt(opts.get("size", 12)); r.font.name = "Consolas"
        r.font.color.rgb = opts.get("color", CODEFG)
    if caption:
        textbox(s, x, y + h + Inches(0.08), w, Inches(0.5), caption,
                size=13, color=GREY)


def image_slide(title, subtitle, img_name, tag=None):
    s = slide()
    title_bar(s, title, subtitle, tag=tag)
    png = os.path.join(BASE, img_name)
    if os.path.exists(png):
        with Image.open(png) as im:
            iw, ih = im.size
        avail_w, avail_h = int(Inches(12.7)), int(Inches(5.9))
        ar = iw / ih
        w = avail_w; h = int(w / ar)
        if h > avail_h:
            h = avail_h; w = int(h * ar)
        x = int((SW - w) / 2); y = int(Inches(1.28)) + int((avail_h - h) / 2)
        s.shapes.add_picture(png, x, y, width=w, height=h)
    else:
        textbox(s, Inches(0.6), Inches(3), Inches(12), Inches(1),
                f"{img_name} not found — run its generator", size=16, color=RED)
    return s


def gap_slide(num, construct, subtitle, legacy_pts, problem_pts, response_pts, footer=None):
    """One detailed slide per legacy gap.
       Left column = what it was (blue) + why it blocks a port (red);
       right light panel = the migration response (green)."""
    s = slide()
    title_bar(s, f"Gap {num} — {construct}", subtitle, tag="gap")
    LW = Inches(6.35)
    # ── left: what it was in the legacy app
    textbox(s, Inches(0.5), Inches(1.26), LW, Inches(0.34),
            "In the legacy app", size=15, color=BLUE, bold=True)
    bullets(s, Inches(0.5), Inches(1.66), LW, Inches(2.5),
            [(t, {**o, "size": o.get("size", 13)}) for t, o in
             ((p if isinstance(p, tuple) else (p, {})) for p in legacy_pts)], size=13, gap=5)
    # ── left: why it blocks a straight port
    textbox(s, Inches(0.5), Inches(4.24), LW, Inches(0.34),
            "Why it blocks a straight port", size=15, color=RED, bold=True)
    bullets(s, Inches(0.5), Inches(4.64), LW, Inches(2.55),
            [(t, {**o, "size": o.get("size", 13), "color": o.get("color", NAVY)}) for t, o in
             ((p if isinstance(p, tuple) else (p, {})) for p in problem_pts)], size=13, gap=5)
    # ── right: migration response (light panel)
    rect(s, Inches(7.05), Inches(1.26), Inches(5.85), Inches(5.95), LIGHT)
    textbox(s, Inches(7.28), Inches(1.4), Inches(5.4), Inches(0.4),
            "Migration response", size=15, color=GREEN, bold=True)
    bullets(s, Inches(7.28), Inches(1.92), Inches(5.4), Inches(4.7),
            [(t, {**o, "size": o.get("size", 13)}) for t, o in
             ((p if isinstance(p, tuple) else (p, {})) for p in response_pts)], size=13, gap=6)
    if footer:
        textbox(s, Inches(7.28), Inches(6.62), Inches(5.4), Inches(0.55),
                footer, size=11, color=GREY)
    return s


def service_slide(name, port, legacy_origin, accent, responsibilities,
                  api_pts, sdk_line, db_line, jms_line, footer=None):
    """One detailed slide per migrated service.
       Left = what it does (responsibilities) + API surface;
       right light panel = the published client SDK, its database, and its JMS role."""
    s = slide()
    title_bar(s, f"{name}  ·  :{port}", legacy_origin, tag="service")
    LW = Inches(6.35)
    # ── left: responsibilities
    textbox(s, Inches(0.5), Inches(1.26), LW, Inches(0.34),
            "What it does", size=15, color=accent, bold=True)
    bullets(s, Inches(0.5), Inches(1.66), LW, Inches(2.7),
            [(t, {**o, "size": o.get("size", 13)}) for t, o in
             ((p if isinstance(p, tuple) else (p, {})) for p in responsibilities)], size=13, gap=5)
    # ── left: API surface
    textbox(s, Inches(0.5), Inches(4.5), LW, Inches(0.34),
            "API surface (HTTP)", size=15, color=BLUE, bold=True)
    bullets(s, Inches(0.5), Inches(4.9), LW, Inches(2.3),
            [(t, {**o, "size": o.get("size", 12)}) for t, o in
             ((p if isinstance(p, tuple) else (p, {})) for p in api_pts)], size=12, gap=4)
    # ── right: SDK / DB / JMS panel
    rect(s, Inches(7.05), Inches(1.26), Inches(5.85), Inches(5.95), LIGHT)
    textbox(s, Inches(7.28), Inches(1.4), Inches(5.4), Inches(0.34),
            "Published client SDK", size=14, color=GREEN, bold=True)
    textbox(s, Inches(7.28), Inches(1.78), Inches(5.4), Inches(1.0),
            [(sdk_line, {"size": 12, "color": NAVY})])
    textbox(s, Inches(7.28), Inches(3.05), Inches(5.4), Inches(0.34),
            "Database", size=14, color=AMBER, bold=True)
    textbox(s, Inches(7.28), Inches(3.43), Inches(5.4), Inches(1.1),
            [(db_line, {"size": 12, "color": NAVY})])
    textbox(s, Inches(7.28), Inches(4.8), Inches(5.4), Inches(0.34),
            "JMS role", size=14, color=RED, bold=True)
    textbox(s, Inches(7.28), Inches(5.18), Inches(5.4), Inches(1.9),
            [(jms_line, {"size": 12, "color": NAVY})])
    if footer:
        textbox(s, Inches(0.5), Inches(7.05), Inches(6.35), Inches(0.4),
                footer, size=11, color=GREY)
    return s


# ═══════════════════════════════════════════════════════════════════════════
# 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.7), SW, Inches(0.06), AMBER)
textbox(s, Inches(0.8), Inches(1.25), Inches(11.7), Inches(1.2),
        "Java Pet Store → Spring Boot 3 / Java 21", size=40, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(2.85), Inches(11.7), Inches(0.8),
        "Migration Playback — gaps fixed · phased plan · code walkthrough", size=24, color=SUB)
textbox(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.2), [
    ("Legacy J2EE 1.3 · 4 EARs · EJB 2.x CMP · JMS (XML) · Cloudscape   →   "
     "8 Spring Boot services · JPA · JMS (JSON) · H2/MongoDB",
     {"size": 15, "color": WHITE, "space_after": 12}),
    ("45-minute session · ~30 min narrative + ~10 min live code + ~5 min Q&A",
     {"size": 14, "color": RGBColor(0x9D, 0xB4, 0xD4)}),
])

# ═══════════════════════════════════════════════════════════════════════════
# 2 — AGENDA (with the 45-min budget)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Agenda", "How the 45 minutes are spent", tag="45 min")
table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), [
    ["#", "Section", "Time"],
    ["1", "The brief & what 'done' means", "2 min"],
    ["2", "Legacy architecture — the starting point", "3 min"],
    ["3", "Gaps in the legacy design (overview + 1 slide per gap)", "8 min"],
    ["4", "Migration strategy & principles", "3 min"],
    ["5", "Phased migration plan (7 phases)", "5 min"],
    ["6", "JMS message migration — XML → JSON (Anti-Corruption Layer)", "4 min"],
    ["7", "What we fixed — parity audit results", "4 min"],
    ["8", "Target as-built architecture", "3 min"],
    ["9", "Service inventory — summary + one slide per service", "8 min"],
    ["10", "CODE WALKTHROUGH (live)", "10 min"],
    ["11", "MongoDB stretch goal", "2 min"],
    ["12", "Wrap-up & Q&A", "5 min"],
], col_widths=[Inches(0.8), Inches(9.5), Inches(2.0)], fsize=13)

# ═══════════════════════════════════════════════════════════════════════════
# 3 — THE BRIEF
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "The Brief — what was asked", "The migration approach is the real deliverable", tag="2 min")
bullets(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), [
    ("Target runtime: latest stable Spring Boot 3.x on Java 21 — host anywhere (laptop / cloud).",
     {"bold": True, "size": 17}),
    ("Public GitHub project with a README covering build + run steps for a developer.",
     {"size": 16}),
    ("Approach it as if the codebase were far larger:", {"bold": True, "size": 17, "color": BLUE}),
    ("break the problem into manageable tasks", {"level": 1}),
    ("apply infrastructure, scaffolding & SW-engineering principles that mitigate migration risk",
     {"level": 1}),
    ("ensure the QUALITY of what is migrated — playback probes this", {"level": 1, "color": RED}),
    ("Optional stretch: run against MongoDB instead of the relational DB.  → DONE.",
     {"bold": True, "size": 17, "color": GREEN}),
    ("Guiding constraint: zero business-logic change — parity first, characterization tests pin behavior.",
     {"size": 16, "color": AMBER, "bold": True}),
])

# ═══════════════════════════════════════════════════════════════════════════
# 4 — LEGACY ARCHITECTURE (embed detailed diagram)
# ═══════════════════════════════════════════════════════════════════════════
image_slide("Legacy Architecture — the starting point",
            "4 EARs · coupled ONLY via JMS · shared Cloudscape DB · messages travel as XML",
            "petstore_legacy_highlevel.png", tag="3 min")

# ═══════════════════════════════════════════════════════════════════════════
# 5 — GAPS IN THE LEGACY DESIGN — overview, then one detailed slide per gap
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Gaps in the Legacy Design — Overview",
          "8 structural gaps block a straight port · each gets its own slide next", tag="gaps")
bullets(s, Inches(0.5), Inches(1.35), Inches(12.4), Inches(5.9), [
    ("These are not bugs — the app worked in 2002. They are design decisions that a modern JVM, "
     "modern tooling, and a service-per-context target render unworkable or unsafe.",
     {"size": 15, "color": NAVY, "space_after": 10}),
    ("Gap 1 — J2EE 1.3 EAR packaging on the Sun Reference Implementation, built with Ant.",
     {"size": 14, "bold": True, "color": BLUE}),
    ("Gap 2 — EJB 2.x: entity beans (CMP), session beans and message-driven beans everywhere.",
     {"size": 14, "bold": True, "color": BLUE}),
    ("Gap 3 — inter-app messages hand-marshalled as XML documents (TPA/DTD) over JMS.",
     {"size": 14, "bold": True, "color": BLUE}),
    ("Gap 4 — one shared Cloudscape database with the same tables duplicated per EAR.",
     {"size": 14, "bold": True, "color": BLUE}),
    ("Gap 5 — order status mutated in three MDBs plus the WAF screen-flow state machine.",
     {"size": 14, "bold": True, "color": BLUE}),
    ("Gap 6 — concurrency safety left implicit to the EJB container (oversell risk).",
     {"size": 14, "bold": True, "color": BLUE}),
    ("Gap 7 — the admin tool is a Swing / Java Web Start rich client.",
     {"size": 14, "bold": True, "color": BLUE}),
    ("Gap 8 — a parallel SOAP 'webservices' build re-implements the same order logic.",
     {"size": 14, "bold": True, "color": BLUE}),
])

# ── Gap 1 — J2EE 1.3 EAR / Sun RI / Ant ──────────────────────────────────────
gap_slide(
    1, "J2EE 1.3 EAR packaging on the Sun Reference Implementation",
    "The runtime itself is gone — before any code, the deployment model must change",
    legacy_pts=[
        ("Ships as 4 enterprise archives (.ear): petstore, opc, supplier, admin — each a bundle "
         "of web (.war) + EJB (.jar) modules glued by application.xml.", {}),
        ("Runs only inside the J2EE 1.3 Sun Reference Implementation app server (Cloudscape DB, "
         "JNDI naming, JMS provider all supplied by that container).", {}),
        ("Built with Ant + hand-written deployment descriptors; deployed by dropping EARs into the "
         "server and editing sun-j2ee-ri.xml bindings.", {}),
        ("Every import is javax.* (javax.ejb, javax.servlet, javax.jms, javax.naming).", {}),
    ],
    problem_pts=[
        ("The Sun RI app server no longer exists and does not run on a modern JVM — there is nothing "
         "to deploy an EAR into on Java 21.", {"color": RED}),
        ("Jakarta EE 9+ renamed every javax.* package to jakarta.* — legacy imports won't even "
         "compile against current libraries.", {"color": RED}),
        ("EAR-in-container assumes a heavyweight app server; it can't be a self-contained, "
         "'java -jar' cloud artifact.", {"color": RED}),
    ],
    response_pts=[
        ("Spring Boot 3.3.5 fat-jars — each service is one self-contained executable jar with an "
         "embedded Tomcat; host on a laptop or any cloud with java -jar.", {"bold": True}),
        ("Maven multi-module build replaces Ant; dependencies + versions are declarative.", {}),
        ("Blanket javax.* → jakarta.* migration across the codebase.", {}),
        ("The 4 EARs become 8 independently deployable services (see the target architecture).", {}),
        ("Config is externalised (application.yml / env vars) instead of container-bound "
         "descriptors — 12-factor friendly.", {}),
    ],
    footer="Outcome: a runnable-anywhere artifact on a supported runtime — the precondition for everything else.")

# ── Gap 2 — EJB 2.x (with the primer folded in) ──────────────────────────────
s = slide()
title_bar(s, "Gap 2 — EJB 2.x is the entire backbone (primer + gap)",
          "Enterprise JavaBeans: the J2EE server-side component model the whole app is built on", tag="gap")
textbox(s, Inches(0.5), Inches(1.24), Inches(12.4), Inches(0.7),
        [("EJB = Enterprise JavaBeans — the standard server-side component model of J2EE. A 'bean' is a "
          "managed object the app server hosts; you write business logic and the container supplies "
          "transactions, security, pooling, remoting, lifecycle and (for entity beans) persistence.",
          {"size": 13, "color": NAVY})])
# left: the three bean types the Pet Store uses
textbox(s, Inches(0.5), Inches(2.15), Inches(6.35), Inches(0.34),
        "The three bean types the Pet Store uses", size=15, color=BLUE, bold=True)
bullets(s, Inches(0.5), Inches(2.55), Inches(6.35), Inches(2.2), [
    ("Session beans — business logic. Stateful ShoppingCart; stateless SignOn, Catalog, "
     "OPC/PurchaseOrder facades.", {"size": 13}),
    ("Entity beans (CMP) — persistent rows the container maps to tables: Customer, Account, Profile, "
     "PurchaseOrder, LineItem, Inventory, ContactInfo, CreditCard.", {"size": 13, "color": AMBER}),
    ("Message-driven beans (MDB) — async JMS consumers: PurchaseOrderMDB, InvoiceMDB, "
     "OrderApprovalMDB, Mail*MDB.", {"size": 13}),
], gap=6)
textbox(s, Inches(0.5), Inches(4.75), Inches(6.35), Inches(0.34),
        "Why it blocks a straight port", size=15, color=RED, bold=True)
bullets(s, Inches(0.5), Inches(5.15), Inches(6.35), Inches(2.1), [
    ("Every bean needs an interface trio (Home + Remote/Local + Bean class) plus ejb-jar.xml / "
     "sun-j2ee-ri.xml wiring — enormous boilerplate.", {"size": 12, "color": NAVY}),
    ("Beans are located by JNDI via a ServiceLocator, not injected — you cannot unit-test without a "
     "running container.", {"size": 12, "color": RED}),
    ("CMP mappings + finder queries (EJB-QL) live in XML, tied to the dead container.", {"size": 12, "color": NAVY}),
], gap=5)
# right: modern mapping
rect(s, Inches(7.05), Inches(2.15), Inches(5.85), Inches(5.05), LIGHT)
textbox(s, Inches(7.28), Inches(2.28), Inches(5.4), Inches(0.4),
        "Migration response — bean-by-bean mapping", size=14, color=GREEN, bold=True)
bullets(s, Inches(7.28), Inches(2.8), Inches(5.4), Inches(4.2), [
    ("Session bean → Spring @Service / @Component (constructor injection, no JNDI).", {"size": 13, "bold": True}),
    ("Stateful ShoppingCart → session-scoped cart-lib component.", {"size": 12, "level": 1}),
    ("Entity bean (CMP) → JPA @Entity + a Spring Data repository behind a domain port.", {"size": 13, "bold": True}),
    ("EJB-QL finders → derived queries / @Query; mappings in annotations, not XML.", {"size": 12, "level": 1}),
    ("Message-driven bean → @JmsListener method on a plain bean.", {"size": 13, "bold": True}),
    ("Home/Remote/Local trios + ServiceLocator → deleted; DI wires collaborators.", {"size": 12, "level": 1}),
    ("Result: plain testable Java — JUnit + Mockito, no container needed.", {"size": 13, "color": GREEN, "bold": True}),
], gap=6)

# ── Gap 3 — XML-over-JMS messages ────────────────────────────────────────────
gap_slide(
    3, "Inter-app messages are hand-marshalled XML over JMS",
    "The apps only talk via JMS — and every payload is a verbose XML document",
    legacy_pts=[
        ("The 4 EARs are decoupled: they never call each other directly, they exchange JMS messages.",
         {}),
        ("Payloads are XML documents built by hand (the TPA — Trading Partner Agreement — schemas, "
         "plus PurchaseOrder / Invoice / SupplierPO DTDs).", {}),
        ("Checkout serialises a PurchaseOrder into XML and sends it to jms/PurchaseOrderQueue; the OPC "
         "MDB parses the XML back into objects.", {}),
        ("Every message type has bespoke marshal / unmarshal code and a matching DTD.", {}),
    ],
    problem_pts=[
        ("Hand-rolled XML marshalling is brittle and verbose — a field rename means editing DTD, "
         "marshaller and parser in lock-step.", {"color": RED}),
        ("DTD-bound payloads are rigid: no forward compatibility, a new field can break older readers.",
         {"color": RED}),
        ("The wire format leaks into domain code — business classes know about XML structure.",
         {"color": RED}),
        ("XML parsing per message is CPU-heavy and hard to debug when a payload is malformed.",
         {"color": RED}),
    ],
    response_pts=[
        ("Model each message as a typed Java record wrapped in one shared event envelope "
         "(petstore-messaging).", {"bold": True}),
        ("Serialize as JSON via MappingJackson2MessageConverter; a _type header routes each JSON back "
         "to the right record.", {}),
        ("An Anti-Corruption Layer maps wire ⇄ domain, so the domain never sees the transport format.",
         {}),
        ("FAIL_ON_UNKNOWN_PROPERTIES=false → a newer producer's additive field can't poison an "
         "older consumer.", {"color": GREEN}),
        ("Golden-file tests freeze the legacy contract before the switch (full detail in the JMS "
         "migration slide, §6).", {}),
    ],
    footer="This is the gap the brief called out explicitly — legacy JMS = XML; migrated JMS = JSON + ACL.")

# ── Gap 4 — shared DB, duplicated tables ─────────────────────────────────────
gap_slide(
    4, "One shared Cloudscape DB, tables duplicated per EAR",
    "No single owner of any entity — the same customer data is copied across apps",
    legacy_pts=[
        ("All 4 EARs read/write one shared Cloudscape database instance.", {}),
        ("Common structures — ContactInfo, Address, CreditCard, LineItem — are redefined and stored "
         "independently inside petstore, opc and supplier.", {}),
        ("There is no authoritative owner: customer/account data physically lives in several schemas "
         "at once.", {}),
        ("CMP entity beans in different EARs map their own copies of overlapping tables.", {}),
    ],
    problem_pts=[
        ("Shared schema = shared coupling: any table change ripples across every EAR at once.",
         {"color": RED}),
        ("Duplicated entities drift — the same 'customer' can disagree between apps, with no source "
         "of truth.", {"color": RED}),
        ("Impossible to scale, deploy or reason about a service in isolation when they all touch one "
         "DB.", {"color": RED}),
    ],
    response_pts=[
        ("Database-per-service: each service owns its schema and is the sole writer of its data.",
         {"bold": True}),
        ("auth-service is the single owner of credentials/accounts; customer-service owns "
         "profile/address/card; OPC owns orders; inventory owns stock.", {}),
        ("Cross-service data is fetched over typed client SDKs (HTTP), never by reaching into another "
         "service's tables.", {}),
        ("Each service uses file H2 by default (OPC + catalog also run on MongoDB) — swappable behind "
         "a repository port.", {}),
        ("A clear ownership boundary makes each service independently testable and deployable.",
         {"color": GREEN}),
    ],
    footer="Bounded contexts replace one shared schema — the core of the service decomposition.")

# ── Gap 5 — order status in 3 MDBs + WAF state machine ───────────────────────
gap_slide(
    5, "Order status mutated in 3 MDBs + the WAF state machine",
    "No single authority for an order's lifecycle — status is written in several places",
    legacy_pts=[
        ("Order state advances as JMS messages ripple through separate MDBs "
         "(PurchaseOrderMDB → OrderApprovalMDB → InvoiceMDB → mailer MDBs).", {}),
        ("Each MDB independently writes the order/purchase-order status it cares about.", {}),
        ("The Web Application Framework (WAF) adds its own screen-flow StateMachine / RequestProcessor "
         "on the web tier.", {}),
        ("The 'current status' of an order is therefore an emergent property of several beans.", {}),
    ],
    problem_pts=[
        ("Status logic is smeared across multiple god-objects — no one class you can point at for "
         "'what are the legal transitions?'.", {"color": RED}),
        ("Concurrent messages can interleave and drive invalid transitions; nothing enforces the "
         "state chart.", {"color": RED}),
        ("The bespoke WAF framework is dead weight — a whole homegrown MVC engine to maintain.",
         {"color": RED}),
    ],
    response_pts=[
        ("One OrderStatusService is the single authority for lifecycle changes.", {"bold": True}),
        ("Status is an explicit enum with guarded transitions (illegal transition → rejected), not "
         "an emergent side effect.", {}),
        ("applyStatusChange runs in one @Transactional unit; a @Version optimistic lock stops an "
         "approve+deny race (→ 409).", {}),
        ("The WAF screen flow is replaced by thin Spring MVC controllers + Thymeleaf.", {}),
        ("Status announcements go out via the transactional outbox — commit and publish can't "
         "diverge (§ walkthrough 5).", {"color": GREEN}),
    ],
    footer="Enum + guards was chosen over the full State pattern — right-sized for this lifecycle.")

# ── Gap 6 — implicit container locking / oversell ────────────────────────────
gap_slide(
    6, "Concurrency safety left implicit to the EJB container",
    "Take the container away and nothing stops two orders overselling the same stock",
    legacy_pts=[
        ("Inventory decrement relied on the EJB container's transaction + entity-bean locking to "
         "serialise concurrent access.", {}),
        ("There is no explicit lock or stock-floor check in the business code — it trusts the "
         "container to prevent lost updates.", {}),
        ("Fulfilment reads stock, decides, then decrements as separate steps within a "
         "container-managed transaction.", {}),
    ],
    problem_pts=[
        ("Once the container is gone, read-then-write is a classic race: two orders both read '5', "
         "both decrement, stock goes negative (oversell).", {"color": RED}),
        ("Nothing in the schema prevents a negative quantity.", {"color": RED}),
        ("With JMS at-least-once delivery, a redelivered fulfilment could double-decrement.",
         {"color": RED}),
    ],
    response_pts=[
        ("Pessimistic row lock: findByIdForUpdate issues SELECT … FOR UPDATE so only one transaction "
         "touches a stock row at a time.", {"bold": True}),
        ("Check-and-decrement runs inside one @Transactional method.", {}),
        ("A DB CHECK (quantity >= 0) constraint is the hard floor — the DB refuses to go negative.",
         {}),
        ("Idempotent consumer: a fulfilled_order dedup ledger means a redelivered order never "
         "double-decrements.", {}),
        ("Verified with a 20-thread test on 5 units → exactly 5 succeed, never negative; the rest "
         "backorder (stay APPROVED).", {"color": GREEN}),
    ],
    footer="This was the one sanctioned technical fix — an explicit lock replacing an implicit one.")

# ── Gap 7 — Swing / Java Web Start admin ─────────────────────────────────────
gap_slide(
    7, "The admin tool is a Swing / Java Web Start rich client",
    "The back-office console is a desktop app launched from the browser — a dead delivery model",
    legacy_pts=[
        ("The admin application is a Swing GUI, delivered via Java Web Start (JNLP) and launched "
         "from a browser link.", {}),
        ("It talks to the OPC server over a remote EJB / XML-RPC channel to list, approve and deny "
         "orders and view sales charts.", {}),
        ("It is a separate rich-client codebase from the web storefront.", {}),
    ],
    problem_pts=[
        ("Java Web Start was removed from the JDK after Java 8 — the launch mechanism no longer "
         "exists.", {"color": RED}),
        ("Browsers dropped the applet/JNLP plumbing entirely — nothing to click.", {"color": RED}),
        ("A desktop Swing client is the wrong shape for a modern web-hosted, deploy-anywhere admin "
         "surface.", {"color": RED}),
    ],
    response_pts=[
        ("Admin becomes REST endpoints on admin-office-service (list / approve / deny / sales) that "
         "delegate to OPC.", {"bold": True}),
        ("A server-rendered Thymeleaf console replaces the Swing UI — same tasks, in the browser.",
         {}),
        ("admin-office-service owns no data — it is a thin console over the OPC facade (SDK).", {}),
        ("Auth is unified: the admin logs in through the same central IdP (JWT) as everything else.",
         {}),
    ],
    footer="Same operator tasks, delivered as a normal web app instead of a desktop download.")

# ── Gap 8 — parallel SOAP webservices build ──────────────────────────────────
gap_slide(
    8, "A parallel SOAP 'webservices' build duplicates order logic",
    "The sample ships two implementations of the supplier flow — JMS and SOAP",
    legacy_pts=[
        ("Pet Store 1.3.1 includes a webservices/ variant that re-implements the supplier purchase-"
         "order exchange over SOAP/JAX-RPC instead of JMS.", {}),
        ("It exists to demo J2EE web-services interoperability — the same business flow, a second "
         "transport.", {}),
        ("Two code paths, two sets of DTOs, for one conceptual operation.", {}),
    ],
    problem_pts=[
        ("Maintaining two implementations of the same order logic doubles the surface with no "
         "functional gain for this project.", {"color": RED}),
        ("JAX-RPC / early SOAP stacks are themselves legacy and unsupported on the target runtime.",
         {"color": RED}),
        ("Keeping both would blur which path is authoritative.", {"color": RED}),
    ],
    response_pts=[
        ("Dropped intentionally: the JMS build is kept as the single source of truth for the "
         "order/supplier flow.", {"bold": True}),
        ("This is a scope decision, recorded as an ADR — not an accidental omission.", {}),
        ("If external SOAP interop were ever needed, it would be re-added as a thin adapter in front "
         "of the same domain, not a parallel implementation.", {}),
    ],
    footer="A deliberate de-duplication — one flow, one implementation.")

# ═══════════════════════════════════════════════════════════════════════════
# 6 — MIGRATION STRATEGY & PRINCIPLES
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Migration Strategy & Principles", "How we de-risk a large migration", tag="3 min")
bullets(s, Inches(0.5), Inches(1.35), Inches(6.4), Inches(5.6), [
    ("Strangler Fig, not big-bang.", {"bold": True, "size": 17, "color": BLUE}),
    ("Migrate one bounded context at a time; keep it runnable at every step.", {"level": 1}),
    ("Characterization tests FIRST.", {"bold": True, "size": 17, "color": BLUE}),
    ("Pin legacy behavior (incl. quirks) before touching it; migrate → green-or-rollback.", {"level": 1}),
    ("Ports & adapters (hexagonal).", {"bold": True, "size": 17, "color": BLUE}),
    ("Domain depends on interfaces; JPA / JMS / Mongo are swappable adapters.", {"level": 1}),
    ("SOLID + ask before choosing a pattern.", {"bold": True, "size": 17, "color": BLUE}),
    ("Enum+guards over State pattern; thin controllers over Command — right-sized.", {"level": 1}),
])
rect(s, Inches(7.15), Inches(1.5), Inches(5.7), Inches(4.9), LIGHT)
textbox(s, Inches(7.4), Inches(1.65), Inches(5.3), Inches(0.5),
        "The safety net (quality gate)", size=16, color=NAVY, bold=True)
bullets(s, Inches(7.4), Inches(2.25), Inches(5.3), Inches(4.0), [
    ("Test pyramid: unit → @DataJpaTest slice → MockMvc → Testcontainers e2e.", {"size": 13}),
    ("Golden-file message tests before any XML→JSON wire change.", {"size": 13, "color": RED}),
    ("Idempotent @JmsListener — JMS is at-least-once.", {"size": 13}),
    ("Parity audit: file-by-file legacy vs migrated (preserved/changed/missing).", {"size": 13}),
    ("Commit only on green; a migration ledger tracks each module.", {"size": 13}),
], gap=9)

# ═══════════════════════════════════════════════════════════════════════════
# 7 — PHASED MIGRATION PLAN
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Phased Migration Plan", "Vertical slices — each phase leaves a runnable app", tag="5 min")
table(s, Inches(0.35), Inches(1.32), Inches(12.65), Inches(5.7), [
    ["Phase", "Scope", "Done when"],
    ["0 · Scaffold", "Parent pom, Boot 3.x/Java 21, H2, Artemis, CI, char-test harness", "empty app boots on Java 21"],
    ["1 · Catalog", "Browse + catalog schema (CatalogRepository port + JPA adapter)", "runs & browsable locally"],
    ["2 · Customer / SignOn", "Register, auth, profile (CMP → JPA); BCrypt upgrade", "auth + profile tests green"],
    ["3 · Cart", "Session cart; all edge-case quirks preserved", "cart edge tests green"],
    ["4 · Order", "Checkout + PurchaseOrder aggregate; OrderStatusService", "checkout persists + emits"],
    ["5 · Fulfilment + JMS", "Inventory (race-fixed), @JmsListener workflow, ACL", "end-to-end order flow green"],
    ["6 · Web + Admin", "Thymeleaf UI, REST admin (replaces Swing), exception handler", "UI + admin flows green"],
    ["7 · Cutover", "Data migration, parallel-run, reconciliation, decommission", "legacy off; ledger all-green"],
], col_widths=[Inches(2.3), Inches(7.55), Inches(2.8)], fsize=11)
textbox(s, Inches(0.35), Inches(7.02), Inches(12.65), Inches(0.4),
        ("Bias to a walking skeleton: Phase 1 gives a browsable app on day one; "
         "every later phase adds one vertical slice without breaking it."),
        size=12, color=GREY)

# ═══════════════════════════════════════════════════════════════════════════
# 8 — JMS MESSAGE MIGRATION: XML → JSON  (explicitly requested)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "JMS Message Migration — XML → JSON", "Legacy sent XML over JMS; we move to a JSON envelope safely", tag="4 min")
textbox(s, Inches(0.5), Inches(1.28), Inches(12.3), Inches(0.55), [
    ("Legacy fact: shared data travelled as hand-marshalled XML documents (TPA/DTD) over JMS — "
     "e.g. checkout built a PurchaseOrder XML → jms/PurchaseOrderQueue.", {"size": 13, "color": RED, "bold": True})])
table(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(2.55), [
    ["Step", "What we do", "Why / risk mitigated"],
    ["1 · Freeze the contract", "Capture real legacy XML messages as golden files (one per destination)",
     "A regression oracle before any change"],
    ["2 · Model events as records", "PurchaseOrder/Invoice/Approved… → typed Java records + a shared envelope",
     "Compile-time contract; no DTD"],
    ["3 · Anti-Corruption Layer", "Map XML DTO ⇄ domain record in messaging/ — domain never sees XML",
     "Wire format isolated behind a port"],
    ["4 · Switch serializer to JSON", "MappingJackson2MessageConverter; _type header for routing",
     "Human-readable, versionable payloads"],
    ["5 · Forward-compat + idempotency", "FAIL_ON_UNKNOWN_PROPERTIES=false; dedup keys on consumers",
     "Additive fields safe; at-least-once safe"],
], col_widths=[Inches(2.7), Inches(5.6), Inches(4.0)], fsize=10.5)
code_panel(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.15),
    "petstore-messaging/.../MessagingConfig.java  (JSON converter + forward-compat)",
    [("var conv = new MappingJackson2MessageConverter();", {}),
     ("conv.setTargetType(MessageType.TEXT);", {}),
     ('conv.setTypeIdPropertyName(\"_type\");   // routes JSON back to the right record', {"color": CODEKW}),
     ("mapper.configure(FAIL_ON_UNKNOWN_PROPERTIES, false); // newer producer can add fields", {"color": CODEKW})])

# ═══════════════════════════════════════════════════════════════════════════
# 9 — WHAT WE FIXED (parity audit)
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "What We Fixed — Parity Audit", "File-by-file legacy vs migrated · 21 gaps found, all resolved", tag="4 min")
table(s, Inches(0.35), Inches(1.32), Inches(12.65), Inches(4.55), [
    ["#", "Gap the audit found", "Resolution"],
    ["H2", "Backorder retry-on-restock dropped", "RESTORED event-driven: RestockEvent → OPC re-drives APPROVED backorders"],
    ["H3/H4", "Approval / denial / COMPLETED emails missing", "OrderStatusTopic → notification-service emails restored"],
    ["H5", "Sales/revenue aggregation gone", "OPC aggregateSales() + GET /api/sales; admin delegates"],
    ["H6", "Catalog search semantics changed", "Restored tokenized multi-field LIKE search"],
    ["H7", "Ship/bill address collection + validation lost", "Both addresses collected + validated; persisted on order"],
    ["H8/H9", "New-customer profile defaults + sign-on locale", "Legacy defaults restored; preferredLanguage applied on login"],
    ["M1–M7", "Ordering, getItem.category, batch approval, etc.", "All restored (see PARITY_AUDIT.md)"],
], col_widths=[Inches(1.1), Inches(6.15), Inches(5.4)], fsize=10.5)
rect(s, Inches(0.35), Inches(6.05), Inches(12.65), Inches(1.05), LIGHT)
textbox(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.9), [
    ("Kept intentional (recorded in DECISIONS.md): H1 all-or-nothing fulfilment · M8 no persisted supplier PO.",
     {"size": 13, "color": NAVY, "bold": True, "space_after": 3}),
    ("Two bugs caught by LIVE end-to-end testing (not unit tests): large orders auto-shipped without approval; "
     "approve→fulfil race (published before commit). Both fixed + regression-tested.",
     {"size": 12, "color": RED})])

# ═══════════════════════════════════════════════════════════════════════════
# 10 — TARGET AS-BUILT ARCHITECTURE (embed)
# ═══════════════════════════════════════════════════════════════════════════
image_slide("Target Architecture — As-Built",
            "8 Spring Boot services · standalone Artemis · 2 queues + 3 topics + DLQ · 5 client SDKs",
            "petstore_architecture.png", tag="3 min")

# ═══════════════════════════════════════════════════════════════════════════
# 10b — SERVICE INVENTORY (summary) — one row per service, then a slide each
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Service Inventory — Summary",
          "8 services · 5 published client SDKs · 2 shared libraries · DB-per-service · one broker",
          tag="summary")
table(s, Inches(0.3), Inches(1.28), Inches(12.75), Inches(4.75), [
    ["Service (port)", "From legacy", "Owns / DB", "Publishes SDK", "JMS role"],
    ["petstore-app-v1 (8080)", "petstore.ear web", "no DB (broker client)", "—", "none (sync REST intake)"],
    ["auth-service (8086)", "SignOn EJB", "accounts · H2", "auth-client", "none"],
    ["customer-service (8081)", "Customer CMP", "profile/addr/card · H2", "customer-service-client", "none"],
    ["catalog-service (8083)", "Catalog EJB", "catalog · H2 / MongoDB", "catalog-service-client", "none"],
    ["order-processing (8088)", "opc.ear", "orders+outbox · H2 / MongoDB", "order-processing-client", "consume+publish (core)"],
    ["inventory-service (8085)", "supplier.ear", "stock+dedup · H2", "inventory-service-client", "consume+publish"],
    ["admin-office (8082)", "admin.ear (Swing)", "no DB (delegates to OPC)", "—", "none"],
    ["notification (8087)", "Mail*MDB", "no DB (stateless)", "—", "consume only (topics+DLQ)"],
], col_widths=[Inches(2.7), Inches(2.05), Inches(2.85), Inches(2.75), Inches(2.4)], fsize=10)
rect(s, Inches(0.3), Inches(6.15), Inches(12.75), Inches(1.0), LIGHT)
textbox(s, Inches(0.5), Inches(6.24), Inches(12.4), Inches(0.85), [
    ("Shared libraries (in-process, not services): cart-lib (session cart) · petstore-messaging "
     "(JMS destinations + JSON event envelope).", {"size": 12, "color": NAVY, "bold": True, "space_after": 3}),
    ("Each stateful service owns its own schema (DB-per-service); services never touch another's tables — "
     "they call its client SDK over HTTP or react to its JMS events. Follow-up slides cover each service.",
     {"size": 12, "color": GREY})])

# ── Service 1 — petstore-app-v1 (storefront) ─────────────────────────────────
service_slide(
    "petstore-app-v1  (Storefront)", 8080, "From the legacy petstore.ear web tier (WAF + JSP)",
    STORE_ACCENT,
    responsibilities=[
        ("The customer-facing shop: browse categories/products/items, search, session cart, checkout.", {}),
        ("Server-rendered Thymeleaf UI with i18n (en/ja/zh) — replaces the legacy WAF + JSP screen flow.", {}),
        ("Persists NOTHING itself: it is a pure client of the backing services + the broker.", {}),
        ("Embeds cart-lib in-process; resolves cart items via the catalog SDK.", {}),
        ("Live stock badge / quantity cap driven by the inventory SDK.", {}),
    ],
    api_pts=[
        ("GET / /category /product /item /search", {}),
        ("POST /cart add/set/update/delete", {}),
        ("GET/POST /checkout  → sync REST intake to OPC (JWT forwarded)", {}),
    ],
    sdk_line="Publishes none (it's an edge app, not a called service). IMPORTS five: "
             "auth-client, catalog-service-client, customer-service-client, "
             "order-processing-client, inventory-service-client.",
    db_line="No database. State lives in the session cart (cart-lib, ~15-min TTL) and in the "
            "downstream services it calls.",
    jms_line="No JMS. Checkout is synchronous REST to OPC (immediate authoritative result). "
             "The legacy fire-and-forget PurchaseOrder-XML publish was replaced by that REST intake "
             "(OPC still keeps a PurchaseOrderQueue listener as an alt async path).",
    footer="Legacy mapping: petstore.ear web tier → thin Spring MVC storefront, no business data of its own.")

# ── Service 2 — auth-service ─────────────────────────────────────────────────
service_slide(
    "auth-service  (central IdP)", 8086, "From the legacy SignOn EJB — now the one identity provider",
    AUTH_ACCENT,
    responsibilities=[
        ("The single identity provider: authenticates users and MINTS RS256 JWTs.", {}),
        ("The ONLY holder of the private signing key and of user credentials.", {}),
        ("Provisions credentials for new customers (called by customer-service).", {}),
        ("Every other service verifies JWTs with the public key only (auth-client) — no shared secret.", {}),
    ],
    api_pts=[
        ("POST /auth/login  → mint JWT", {}),
        ("POST /auth/accounts  → provision credential", {}),
    ],
    sdk_line="Publishes auth-client — RS256 verify (public key) + a login helper. Imported by "
             "storefront, customer, admin-office, OPC, inventory (everyone that authenticates or "
             "verifies a token).",
    db_line="Owns the account/credential store · file H2. Sole owner of identity data — "
            "no other service stores credentials.",
    jms_line="None. Auth is request/response only.",
    footer="Legacy mapping: SignOn EJB → a standalone IdP; centralises what was scattered auth logic.")

# ── Service 3 — customer-service ─────────────────────────────────────────────
service_slide(
    "customer-service", 8081, "From the legacy Customer / Account CMP entity beans",
    CUST_ACCENT,
    responsibilities=[
        ("Owns customer PII: profile, contact/address, and credit-card data.", {}),
        ("Registers new customers, then calls auth-service to provision their credential.", {}),
        ("Restores legacy new-customer profile defaults + preferred-language handling.", {}),
    ],
    api_pts=[
        ("POST /register", {}),
        ("GET /customer/{id}", {}),
        ("PUT /account /profile /card", {}),
    ],
    sdk_line="Publishes customer-service-client — customer API + DTOs. Imported by the storefront "
             "(profile/registration screens).",
    db_line="Owns the customer schema (profile / address / card) · file H2. The single owner of "
            "customer PII — replaces the CMP tables duplicated across legacy EARs.",
    jms_line="None. Imports auth-client to provision credentials over HTTP.",
    footer="Legacy mapping: Customer/Account CMP beans → a bounded customer context with its own schema.")

# ── Service 4 — catalog-service ──────────────────────────────────────────────
service_slide(
    "catalog-service", 8083, "From the legacy Catalog EJB — read-mostly, multi-locale",
    CAT_ACCENT,
    responsibilities=[
        ("Serves the product catalog: categories, products, items, keyword search.", {}),
        ("Read-mostly, multi-locale (locale-split tables/collections; en/ja/zh).", {}),
        ("Restored the legacy tokenized multi-field LIKE search semantics (parity H6).", {}),
        ("Second service (with OPC) proven on MongoDB behind the same repository port.", {}),
    ],
    api_pts=[
        ("GET /api/categories /products /items", {}),
        ("GET /api/items?keyword=...", {}),
    ],
    sdk_line="Publishes catalog-service-client — catalog API + DTOs. Imported by the storefront and, "
             "in-process, by cart-lib (to resolve cart line items).",
    db_line="Owns the catalog schema · file H2 by default, OR MongoDB (@Profile mongo) — 3 "
            "collections with an embedded locale map. Same port, swappable adapter.",
    jms_line="None. Pure read API.",
    footer="Legacy mapping: Catalog EJB → a read-optimised service; the ItemDto still carries the legacy attr1..5 slots.")

# ── Service 5 — order-processing (OPC) ───────────────────────────────────────
service_slide(
    "order-processing-service  (OPC)", 8088, "From the legacy opc.ear — the authoritative order owner",
    OPC_ACCENT,
    responsibilities=[
        ("The authoritative owner of orders and their workflow status for the whole system.", {}),
        ("Intake → auto-approve under a currency threshold, else PENDING for a human admin.", {}),
        ("Owns the order lifecycle enum + guarded transitions; @Version stops approve/deny races.", {}),
        ("Transactional outbox: order write + event emit commit atomically (no lost/dup publish).", {}),
        ("Re-drives APPROVED backorders on RestockTopic (the restored legacy processPendingPO).", {}),
    ],
    api_pts=[
        ("POST /api/orders/intake  (sync checkout)", {}),
        ("GET/POST /api/orders  approve|deny|status", {}),
        ("GET /api/sales  (revenue aggregation)", {}),
    ],
    sdk_line="Publishes order-processing-client — OPC facade + intake DTOs + endpoint constants. "
             "Imported by the storefront (checkout intake) and admin-office (approve/deny/sales).",
    db_line="Owns orders + outbox (authoritative) · file H2 by default (Flyway schema), OR MongoDB "
            "(@Profile mongo) with multi-document transactions on replica set rs0.",
    jms_line="The JMS hub. CONSUMES PurchaseOrderQueue, InvoiceTopic, RestockTopic; "
             "PUBLISHES OrderApprovedEvent → ApprovedOrderQueue and OrderStatusEvent → "
             "OrderStatusTopic. All outbound events go through the outbox relay.",
    footer="Legacy mapping: opc.ear (PurchaseOrderMDB + OPCAdminFacade) → one service, single order writer.")

# ── Service 6 — inventory-service ────────────────────────────────────────────
service_slide(
    "inventory-service  (fulfilment)", 8085, "From the legacy supplier.ear — fulfils approved orders",
    INV_ACCENT,
    responsibilities=[
        ("Fulfils approved orders: reserves stock and ships, or backorders if short.", {}),
        ("Oversell guard: pessimistic SELECT…FOR UPDATE + DB CHECK(qty>=0) (the one sanctioned fix).", {}),
        ("Idempotent: a fulfilled_order dedup ledger stops a redelivered order double-decrementing.", {}),
        ("Suppliers restock stock, which re-drives backorders across the system.", {}),
    ],
    api_pts=[
        ("GET /api/inventory  (stock read)", {}),
        ("POST /api/inventory  restock (SUPPLIER role)", {}),
    ],
    sdk_line="Publishes inventory-service-client — stock read SDK with a single-flight cache. "
             "Imported by the storefront (live stock badge / quantity cap). The newest SDK.",
    db_line="Owns the inventory schema (on-hand stock) + the dedup ledger · file H2. Sole owner "
            "of stock levels.",
    jms_line="CONSUMES ApprovedOrderQueue (reserve+ship). PUBLISHES InvoiceEvent → InvoiceTopic "
             "(ship confirmation) and RestockEvent → RestockTopic (triggers OPC's backorder re-drive).",
    footer="Legacy mapping: supplier.ear → a fulfilment service; container locking replaced by an explicit row lock.")

# ── Service 7 — admin-office-service ─────────────────────────────────────────
service_slide(
    "admin-office-service  (console)", 8082, "From the legacy admin.ear Swing / Java Web Start client",
    ADMIN_ACCENT,
    responsibilities=[
        ("The back-office console: list / approve / deny orders, view sales.", {}),
        ("Server-rendered Thymeleaf UI replacing the dead Swing / Java Web Start rich client.", {}),
        ("Owns NO data — a thin console that delegates every action to OPC via its SDK.", {}),
        ("Admin logs in through the same central IdP (JWT) as everyone else.", {}),
    ],
    api_pts=[
        ("GET /warehouse/orders  (console UI)", {}),
        ("POST approve | deny", {}),
    ],
    sdk_line="Publishes none. IMPORTS order-processing-client (all order/sales actions) and "
             "auth-client (login + verify).",
    db_line="No database. Every read/write is delegated to OPC — admin-office is stateless.",
    jms_line="None. Pure HTTP delegation to OPC.",
    footer="Legacy mapping: admin.ear Swing client → a normal web console; same operator tasks in the browser.")

# ── Service 8 — notification-service ─────────────────────────────────────────
service_slide(
    "notification-service", 8087, "From the legacy Mail*MDB beans — a pure event observer",
    NOTE_ACCENT,
    responsibilities=[
        ("Sends customer emails: approval / denial / shipped / completed.", {}),
        ("A pure, stateless event observer — persists nothing, owns no business state.", {}),
        ("Also the operator safety net: logs DLQ + ExpiryQueue messages at ERROR.", {}),
    ],
    api_pts=[
        ("No public HTTP API — it is a message consumer only.", {}),
    ],
    sdk_line="Publishes none — nothing calls it. It only subscribes to the broker.",
    db_line="No database — stateless. (Idempotency is deferred by decision: harmless while the "
            "mailer only logs.)",
    jms_line="CONSUME-ONLY. Subscribes InvoiceTopic + OrderStatusTopic (emails) and the "
             "DLQ + ExpiryQueue (ERROR-log observer). Publishes nothing.",
    footer="Legacy mapping: the Mail*MDB beans → one notification service; also closes the DLQ 'black hole'.")

# ═══════════════════════════════════════════════════════════════════════════
# 11 — CODE WALKTHROUGH divider
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.5), SW, Inches(0.06), AMBER)
textbox(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.0),
        "Code Walkthrough", size=40, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(3.75), Inches(11.7), Inches(1.6), [
    ("Live in the repo — the slides mirror the files, open the real code alongside.", {"size": 17, "color": SUB, "space_after": 10}),
    ("1) JMS contract  2) Anti-Corruption / envelope  3) checkout intake  "
     "4) pessimistic-lock fulfilment  5) restock re-drive  6) transactional outbox", {"size": 14, "color": RGBColor(0x9D, 0xB4, 0xD4)}),
], anchor=MSO_ANCHOR.TOP)

# ── 12 — WALKTHROUGH 1: the JMS contract (one shared lib) ─────────────────────
s = slide()
title_bar(s, "Walkthrough 1 — The JMS Contract", "One shared library defines every destination; producers/consumers import it", tag="code")
code_panel(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.55),
    "petstore-messaging/src/com/petstore/messaging/Destinations.java",
    [('public static final String PURCHASE_ORDER_NAME = \"PurchaseOrderQueue\";', {}),
     ('public static final String APPROVED_ORDER_NAME = \"ApprovedOrderQueue\";', {}),
     ('public static final String INVOICE_NAME        = \"InvoiceTopic\";', {}),
     ('public static final String ORDER_STATUS_NAME   = \"OrderStatusTopic\";', {}),
     ('public static final String RESTOCK_NAME        = \"RestockTopic\";', {}),
     ("", {}),
     ("public static final Destination APPROVED_ORDER = queue(APPROVED_ORDER_NAME);", {"color": CODEKW}),
     ("public static final Destination INVOICE        = topic(INVOICE_NAME);   // pub/sub fan-out", {"color": CODEKW}),
     ("public static final Destination RESTOCK        = topic(RESTOCK_NAME);   // H2 re-drive", {"color": CODEKW})])
bullets(s, Inches(0.5), Inches(5.15), Inches(12.3), Inches(2.0), [
    ("Single source of truth — no service hard-codes a destination string.", {"size": 14}),
    ("queue() = point-to-point (one consumer) · topic() = fan-out (legacy InvoiceTopic preserved).", {"size": 14}),
    ("2 queues + 3 topics; DLQ + ExpiryQueue are broker-managed safety nets.", {"size": 14}),
])

# ── 13 — WALKTHROUGH 2: synchronous checkout intake ───────────────────────────
s = slide()
title_bar(s, "Walkthrough 2 — Checkout Intake (sync REST)", "Storefront hands the order to OPC over REST; JWT proxied for authz", tag="code")
code_panel(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.0),
    "petstore-app-v1/src/com/petstore/order/service/OrderService.java",
    [("List<CartItem> items = cart.getItems();", {}),
     ("if (items.isEmpty()) throw new EmptyCartException();", {}),
     ("// LOCALE = Locale.US, CURRENCY = \"USD\" (legacy quirk — UI locale does not flow in)", {"color": GREEN}),
     ("CheckoutRequest request = new CheckoutRequest(", {}),
     ("    orderId, userId, emailId, LOCALE, CURRENCY, total, lines, toDto(shipTo), toDto(billTo));", {}),
     ("try {", {}),
     ("    response = orderProcessing.checkout(request, bearer); // POST /api/orders/intake, JWT fwd", {"color": CODEKW}),
     ("} catch (RestClientException e) {", {}),
     ("    throw new OrderIntakeUnavailableException(...); // OPC down → 503, cart NOT emptied", {"color": RED}),
     ("}", {}),
     ("cart.empty();   // only on success", {"color": CODEKW})])
bullets(s, Inches(0.5), Inches(4.65), Inches(12.3), Inches(2.35), [
    ("Legacy published a PurchaseOrder XML fire-and-forget; we made intake synchronous so the shopper "
     "gets an immediate, authoritative result.", {"size": 14}),
    ("The storefront persists nothing — OPC is the authoritative order store (single writer).", {"size": 14}),
    ("PurchaseOrderQueue listener is kept as an alternate async intake path (contract unchanged).", {"size": 14, "color": GREY}),
])

# ── 14 — WALKTHROUGH 3: pessimistic-lock fulfilment (oversell guard) ──────────
s = slide()
title_bar(s, "Walkthrough 3 — Oversell Guard", "The sanctioned technical fix: pessimistic row lock replaces the EJB container lock", tag="code")
code_panel(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.85),
    "inventory-service/.../repository/jpa/InventoryJpaRepository.java",
    [("@Lock(LockModeType.PESSIMISTIC_WRITE)", {"color": CODEKW}),
     ('@Query(\"select i from InventoryEntity i where i.itemId = :itemId\")', {}),
     ("Optional<InventoryEntity> findByIdForUpdate(String itemId);  // SELECT … FOR UPDATE", {}),
     ("", {}),
     ("@Modifying(flushAutomatically = true, clearAutomatically = true)", {"color": CODEKW}),
     ('@Query(\"update InventoryEntity i set i.quantity = i.quantity + :qty where i.itemId = :itemId\")', {}),
     ("int increment(String itemId, int qty);   // additive restock")])
bullets(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.4), [
    ("check-and-decrement inside one @Transactional; DB CHECK (quantity >= 0) is the floor.", {"size": 14}),
    ("Verified with a 20-thread test: 5 stock → exactly 5 succeed, never negative.", {"size": 14, "color": GREEN}),
    ("Preserves the legacy observable outcome: one order ships, the other backorders (stays APPROVED).", {"size": 14}),
])

# ── 15 — WALKTHROUGH 4: restock re-drive (H2, event-driven) ──────────────────
s = slide()
title_bar(s, "Walkthrough 4 — Backorder Re-drive (H2)", "Restore legacy processPendingPO without the legacy PENDING PO store", tag="code")
code_panel(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.35),
    "order-processing-service/app/.../messaging/RestockListener.java",
    [("@JmsListener(destination = Destinations.RESTOCK_NAME, containerFactory = \"topicFactory\",", {"color": CODEKW}),
     ('           subscription = \"opc-restock\")   // durable shared subscription', {"color": CODEKW}),
     ("@Transactional", {"color": CODEKW}),
     ("public void onRestock(RestockEvent restock) {", {}),
     ("    Correlation.set(restock.meta() == null ? null : restock.meta().correlationId());", {}),
     ("    // AdminService: orderIdsByStatus(APPROVED) → sorted by created (oldest-first)", {"color": GREEN}),
     ("    //             → approvalGateway.dispatchForFulfilment(order) per order", {"color": GREEN}),
     ("    admin.redriveApprovedForFulfilment();", {}),
     ("}", {})])
bullets(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), [
    ("inventory publishes RestockEvent → RestockTopic when stock is added; OPC owns the order read-model and reacts.",
     {"size": 13}),
    ("Idempotent end-to-end: inventory dedups by order_id (fulfilled_order ledger) — a re-driven order that "
     "already shipped never double-decrements.", {"size": 13, "color": GREEN}),
    ("Same observable behavior as legacy; no persisted supplier PO required (M8 kept).", {"size": 13, "color": GREY}),
])

# ── 16 — WALKTHROUGH 5: transactional outbox ─────────────────────────────────
s = slide()
title_bar(s, "Walkthrough 5 — Transactional Outbox", "Order commit and 'approved' publish can't diverge", tag="code")
bullets(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.0), [
    ("Problem: if we persist the order then publish, a crash between them loses the message (or double-sends).",
     {"size": 15, "color": RED}),
    ("Fix: write the outgoing event into an outbox row in the SAME DB transaction as the order; a relay "
     "polls and publishes, marking rows sent. Publish is now exactly-effectively-once.", {"size": 15}),
])
code_panel(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(2.4),
    "order-processing-service/app/.../service/  (AdminService · ApprovalGateway · OutboxRelay)",
    [("@Transactional                                     // ONE atomic unit", {"color": CODEKW}),
     ("void applyStatusChange(orderId, APPROVED) {", {}),
     ("    orders.updateStatus(orderId, APPROVED);         // order row", {}),
     ("    approvalGateway.dispatchForFulfilment(order);   // → OutboxWriter.enqueue(...) — SAME tx", {"color": CODEKW}),
     ("    statusGateway.announce(order, APPROVED);        // → OutboxWriter.enqueue(...) — SAME tx", {"color": CODEKW}),
     ("}", {}),
     ("// @Scheduled OutboxRelay polls unsent rows → MessagePublisher → marks published", {"color": GREEN})])
textbox(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.95),
        ("Gateways NEVER publish to JMS directly — they append to the outbox table in-transaction, so a rolled-back "
         "order emits nothing and a committed one always eventually does. At-least-once relay is safe: fixed eventId "
         "+ idempotent consumers. @Version optimistic lock stops an approve+deny race (→ 409); the Mongo adapter keeps "
         "outbox + @Version via multi-document transactions (rs0)."),
        size=12, color=GREY)

# ═══════════════════════════════════════════════════════════════════════════
# 17 — MONGODB STRETCH
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "MongoDB Stretch — DONE", "Ports made the DB swap additive — zero change to domain / services / tests", tag="2 min")
bullets(s, Inches(0.5), Inches(1.4), Inches(6.5), Inches(5.4), [
    ("Profile-selectable: @Profile(\"mongo\") vs @Profile(\"!mongo\"); H2 stays the default.",
     {"size": 15, "bold": True}),
    ("order-processing-service — fully wired:", {"size": 15, "color": BLUE, "bold": True}),
    ("MongoOrderStore + MongoOutboxStore adapters", {"level": 1, "size": 14}),
    ("multi-document transactions (outbox + order commit)", {"level": 1, "size": 14}),
    ("@Version optimistic lock + $jsonSchema validators", {"level": 1, "size": 14}),
    ("catalog-service — Mongo adapter (3 collections, locale map embedded).", {"size": 15, "color": BLUE, "bold": True}),
    ("Testcontainers mongo:7 parity tests on both.", {"size": 14, "color": GREEN}),
])
rect(s, Inches(7.25), Inches(1.5), Inches(5.6), Inches(4.6), LIGHT)
textbox(s, Inches(7.5), Inches(1.65), Inches(5.1), Inches(0.5),
        "Why it was cheap", size=16, color=NAVY, bold=True)
bullets(s, Inches(7.5), Inches(2.25), Inches(5.1), Inches(3.7), [
    ("The domain talks to a port (OrderStore / OutboxStore), not to JPA.", {"size": 14}),
    ("Mongo is just another adapter behind the same interface.", {"size": 14}),
    ("Run it: OPC_STORE=mongo ./run-all.sh (docker-compose brings up mongo + mongo-express).", {"size": 13, "color": GREY}),
    ("This is the payoff of hexagonal boundaries chosen up front.", {"size": 14, "color": GREEN, "bold": True}),
])

# ═══════════════════════════════════════════════════════════════════════════
# 18 — WRAP-UP
# ═══════════════════════════════════════════════════════════════════════════
s = slide()
title_bar(s, "Wrap-up", "What the playback should take away", tag="5 min")
bullets(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), [
    ("The approach de-risked a large migration: Strangler Fig + characterization tests + ports/adapters.",
     {"size": 17, "bold": True, "color": BLUE}),
    ("Every legacy construct has a deliberate modern mapping — nothing dropped by accident.", {"size": 16}),
    ("Messages moved XML → JSON behind an Anti-Corruption Layer, guarded by golden-file tests.", {"size": 16}),
    ("A file-by-file parity audit found 21 drifts; all resolved or explicitly kept with an ADR.", {"size": 16}),
    ("Live e2e testing caught two bugs unit tests missed — the safety net earned its keep.", {"size": 16, "color": RED}),
    ("Stretch goal delivered: MongoDB behind the same ports, at near-zero domain cost.", {"size": 16, "color": GREEN}),
    ("", {"size": 8}),
    ("Questions?", {"size": 22, "bold": True, "color": NAVY}),
])

# ═══════════════════════════════════════════════════════════════════════════
out = os.path.join(BASE, "PetStore_Demo_45min.pptx")
prs.save(out)
print("Wrote", out)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
