#!/usr/bin/env python3
"""
Pet Store — AS-BUILT high-level architecture (matches docs/ARCHITECTURE.md).

Renders two clean, presentation-grade diagrams with Graphviz, then appends two
slides to the existing deck (docs/PetStore_Architecture_LLD.pptx) NON-DESTRUCTIVELY
(existing slides are preserved; two new ones are added at the end).

Outputs:
  docs/petstore_asbuilt_services.png / .svg   — service map (clients imported + DB type)
  docs/petstore_asbuilt_jms.png     / .svg    — JMS topology (queues/topics + consumers)
  docs/PetStore_Architecture_LLD.pptx          — +2 slides appended

Run:  python3 docs/petstore_architecture_asbuilt.py
Deps: graphviz (+ system `dot`), python-pptx.
"""
import os
import graphviz
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

BASE = os.path.dirname(os.path.abspath(__file__))

# ---- Palette (shared with the deck: petstore_ppt.py) --------------------------
NAVY  = "#1F3355"; BLUE = "#2E6DB4"; AMBER = "#C97C2F"; GREEN = "#2F8F46"
RED   = "#B23A48"; PURPLE = "#6B4FA0"; GREY = "#555555"; TEAL = "#2A9D8F"
LIGHT = "#F2F4F7"; WHITE = "#FFFFFF"

# per-service (fill, border) — kept legible with white title bars
C = {
    "store":  ("#DCE9F7", BLUE),    # petstore-app-v1
    "auth":   ("#E9E1F5", PURPLE),  # auth-service
    "catalog":("#D8F0DD", GREEN),   # catalog-service
    "customer":("#FBE7D0", AMBER),  # customer-service
    "opc":    ("#F7D9DC", RED),     # order-processing
    "admin":  ("#EDEDED", GREY),    # admin-office
    "inv":    ("#D6EEF0", TEAL),    # inventory
    "notif":  ("#FFF3CD", "#B8860B"),  # notification
}
C_DB = ("#FFF3CD", "#B8860B")


def _svc_node(g, nid, title, port, clients, db, extra=""):
    """A service card: title bar, imported-clients line, and DB/persistence line."""
    fill, border = C[nid]
    cl = "clients: [" + ", ".join(clients) + "]" if clients else "clients: [ — ]"
    db_line = f"DB: {db}"
    rows = (
        f'<TR><TD BGCOLOR="{border}"><FONT COLOR="white" POINT-SIZE="14"><B>{title}</B>'
        f'</FONT><FONT COLOR="white" POINT-SIZE="10">  :{port}</FONT></TD></TR>'
        f'<TR><TD BGCOLOR="{fill}" ALIGN="LEFT"><FONT POINT-SIZE="10" FACE="Courier">{cl}</FONT></TD></TR>'
    )
    if extra:
        rows += (f'<TR><TD BGCOLOR="{fill}" ALIGN="LEFT"><FONT POINT-SIZE="9"><I>{extra}</I></FONT></TD></TR>')
    dbcolor = C_DB[1] if db.startswith("SQL") else "#999999"
    rows += (f'<TR><TD BGCOLOR="{fill}"><FONT POINT-SIZE="10" COLOR="{dbcolor}"><B>{db_line}</B></FONT></TD></TR>')
    g.node(nid, f'<<TABLE BORDER="0" CELLBORDER="1" CELLSPACING="0" CELLPADDING="7">{rows}</TABLE>>',
           shape="plaintext")


def _db_node(g, nid, label):
    g.node(nid, label, shape="cylinder", style="filled", fillcolor=C_DB[0],
           color=C_DB[1], fontsize="9", fontcolor=C_DB[1], fontname="Helvetica")


# ============================================================================
# Diagram 1 — SERVICE MAP: services, clients imported, DB type
# ============================================================================
def build_services():
    g = graphviz.Digraph("asbuilt_svc", format="png")
    g.attr(rankdir="TB", splines="spline", nodesep="0.5", ranksep="0.85", bgcolor="white",
           fontname="Helvetica", fontsize="16", labelloc="t", pad="0.3",
           label="Java Pet Store — As-Built Service Map (Spring Boot 3.3.5 / Java 21)\\l"
                 "Service[client SDKs imported]  ·  DB: file H2 (opc/catalog also MongoDB-capable) or none  ·  "
                 "solid blue = HTTP/REST call  ·  dashed red = JMS\\l")
    g.attr("node", fontname="Helvetica")
    g.attr("edge", fontname="Helvetica")

    # actors
    g.node("browser", "Browser\\n(shopper · admin · supplier)", shape="box",
           style="filled,rounded", fillcolor="#ffffff", color=NAVY, fontsize="12", fontcolor=NAVY)

    # edge-facing tier
    _svc_node(g, "store", "petstore-app-v1 (Storefront)", 8080,
              ["auth", "catalog", "customer", "opc", "inventory"], "none (broker client)",
              "embeds cart-lib · sync REST checkout intake → opc · live stock badge")
    _svc_node(g, "admin", "admin-office-service", 8082,
              ["auth", "opc"], "none (delegates)", "back-office ADMIN console")
    _svc_node(g, "inv", "inventory-service", 8085,
              ["auth"], "SQL (h2 inventory)",
              "publishes inventory-service-client · stock + dedup ledger")

    # backing services
    _svc_node(g, "auth", "auth-service (IdP)", 8086, [], "SQL (h2 auth)", "mints RS256 JWT (private key)")
    _svc_node(g, "catalog", "catalog-service", 8083, [], "SQL (h2) or MongoDB", "read-only · locale-split")
    _svc_node(g, "customer", "customer-service", 8081, ["auth"], "SQL (h2 customer)", "PII / profile / card")
    _svc_node(g, "opc", "order-processing (OPC)", 8088, ["auth"], "SQL (h2) or MongoDB",
              "authoritative order store + outbox · uses petstore-messaging")
    _svc_node(g, "notif", "notification-service", 8087, [], "none (stateless observer)", "uses petstore-messaging")

    # rank hints so the layout reads top→bottom in tiers
    with g.subgraph() as s:
        s.attr(rank="same"); s.node("store"); s.node("admin"); s.node("inv")
    with g.subgraph() as s:
        s.attr(rank="same"); s.node("catalog"); s.node("customer"); s.node("opc")
    with g.subgraph() as s:
        s.attr(rank="same"); s.node("auth"); s.node("notif")

    # browser → edge tier
    g.attr("edge", color=NAVY, penwidth="1.6", fontsize="9", fontcolor=NAVY, style="solid")
    g.edge("browser", "store", label=":8080")
    g.edge("browser", "admin", label=":8082")
    g.edge("browser", "inv",   label=":8085")

    # HTTP/REST client calls (solid blue)
    g.attr("edge", color=BLUE, penwidth="1.5", fontsize="9", fontcolor=BLUE, style="solid")
    g.edge("store", "auth",     label="login")
    g.edge("store", "catalog",  label="browse")
    g.edge("store", "customer", label="profile")
    g.edge("store", "opc",      label="checkout intake")
    g.edge("store", "inv",      label="stock badge", constraint="false")
    g.edge("customer", "auth",  label="provision")
    g.edge("admin", "opc",      label="orders / approve / sales")
    g.edge("admin", "auth",     label="login", constraint="false")
    g.edge("opc", "auth",       label="verify", style="dotted", constraint="false")
    g.edge("inv", "auth",       label="verify", style="dotted", constraint="false")

    # databases (one per stateful service)
    g.attr("edge", color=C_DB[1], penwidth="1.2", fontsize="8", fontcolor=C_DB[1], style="solid")
    for svc, dbid, lbl, store in [("auth", "db_auth", "auth", "file H2"),
                                  ("catalog", "db_cat", "catalog", "H2 / MongoDB"),
                                  ("customer", "db_cust", "customer", "file H2"),
                                  ("opc", "db_opc", "opc", "H2 / MongoDB"),
                                  ("inv", "db_inv", "inventory", "file H2")]:
        _db_node(g, dbid, f"{lbl} db\\n({store})")
        g.edge(svc, dbid)

    out = os.path.join(BASE, "petstore_asbuilt_services")
    g.render(out, format="png", cleanup=True)
    g.render(out, format="svg", cleanup=True)
    print("wrote", out + ".png / .svg")
    return out + ".png"


# ============================================================================
# Diagram 2 — JMS TOPOLOGY: one broker, queues/topics, producers & consumers
# ============================================================================
def build_jms():
    g = graphviz.Digraph("asbuilt_jms", format="png")
    g.attr(rankdir="LR", splines="spline", nodesep="0.45", ranksep="1.4", bgcolor="white",
           fontname="Helvetica", fontsize="16", labelloc="t", pad="0.3",
           label="Java Pet Store — JMS Topology (one ActiveMQ Artemis broker, :61616)\\l"
                 "▣ queue = point-to-point (one consumer)   ◈ topic = pub/sub (fan-out)   "
                 "retry ~1s→~2s→3rd fail ⇒ DLQ\\l"
                 "checkout intake is now synchronous REST → opc (PurchaseOrderQueue kept as alt path)\\l")
    g.attr("node", fontname="Helvetica")
    g.attr("edge", fontname="Helvetica")

    def producer(nid, label, col):
        g.node(nid, label, shape="box", style="filled,rounded", fillcolor=C[col][0],
               color=C[col][1], fontsize="11", fontcolor=NAVY)

    def dest(nid, label, kind):
        # queue = amber box, topic = green folder-ish, safety = red
        if kind == "queue":
            fill, border, glyph = "#FBE7D0", AMBER, "▣ QUEUE"
        elif kind == "topic":
            fill, border, glyph = "#D8F0DD", GREEN, "◈ TOPIC"
        else:
            fill, border, glyph = "#F7D9DC", RED, "▣ SAFETY"
        g.node(nid, f'<<TABLE BORDER="0" CELLBORDER="1" CELLSPACING="0" CELLPADDING="6">'
                    f'<TR><TD BGCOLOR="{border}"><FONT COLOR="white" POINT-SIZE="9"><B>{glyph}</B></FONT></TD></TR>'
                    f'<TR><TD BGCOLOR="{fill}"><FONT POINT-SIZE="11"><B>{label}</B></FONT></TD></TR>'
                    f'</TABLE>>', shape="plaintext")

    # producers (left)
    producer("p_opc",   "order-processing\\nOutboxRelay", "opc")
    producer("p_inv",   "inventory-service\\nOrderApproved + Restock", "inv")
    producer("brk",     "broker\\n(after 3 failed deliveries)", "admin")

    # destinations (center)
    dest("q_po",   "PurchaseOrderQueue", "queue")
    dest("q_ao",   "ApprovedOrderQueue", "queue")
    dest("t_inv",  "InvoiceTopic", "topic")
    dest("t_os",   "OrderStatusTopic", "topic")
    dest("t_rs",   "RestockTopic", "topic")
    dest("dlq",    "DLQ", "safety")
    dest("exp",    "ExpiryQueue", "safety")

    # consumers (right)
    def consumer(nid, label, col):
        g.node(nid, label, shape="box", style="filled,rounded", fillcolor=C[col][0],
               color=C[col][1], fontsize="11", fontcolor=NAVY)
    consumer("c_opc",   "order-processing\\nOrder/Invoice/Restock listeners", "opc")
    consumer("c_inv",   "inventory-service\\nOrderApprovedListener", "inv")
    consumer("c_notif", "notification-service\\nInvoice/Status/Dlq listeners", "notif")

    # produce edges (solid) → destination
    g.attr("edge", color=NAVY, penwidth="1.5", fontsize="9", fontcolor=GREY, style="solid", arrowhead="normal")
    g.edge("p_opc",   "q_ao", label="OrderApprovedEvent")
    g.edge("p_opc",   "t_os", label="OrderStatusEvent")
    g.edge("p_inv",   "t_inv", label="InvoiceEvent")
    g.edge("p_inv",   "t_rs",  label="RestockEvent")
    g.edge("brk",     "dlq", label="quarantine", color=RED, fontcolor=RED)
    g.edge("brk",     "exp", label="expired", color=RED, fontcolor=RED)

    # destination → consumer (dashed red = JMS delivery)
    g.attr("edge", color=RED, penwidth="1.8", fontsize="9", fontcolor=RED, style="dashed", arrowhead="vee")
    g.edge("q_po",  "c_opc",   label="(alt intake)")
    g.edge("q_ao",  "c_inv")
    g.edge("t_inv", "c_opc",   label="→ COMPLETED")
    g.edge("t_inv", "c_notif", label="→ email")   # topic fan-out: two consumers
    g.edge("t_os",  "c_notif", label="→ email")
    g.edge("t_rs",  "c_opc",   label="→ re-drive backorders")
    g.edge("dlq",   "c_notif", label="ERROR log")
    g.edge("exp",   "c_notif", label="ERROR log")

    out = os.path.join(BASE, "petstore_asbuilt_jms")
    g.render(out, format="png", cleanup=True)
    g.render(out, format="svg", cleanup=True)
    print("wrote", out + ".png / .svg")
    return out + ".png"


# ============================================================================
# PPT — append two slides to the existing deck (non-destructive)
# ============================================================================
def _fit(img_path, avail_w_emu, avail_h_emu):
    """Return (w, h) in EMU that fits the image inside the available box, keeping aspect."""
    with Image.open(img_path) as im:
        iw, ih = im.size
    ar = iw / ih
    w = avail_w_emu
    h = int(w / ar)
    if h > avail_h_emu:
        h = avail_h_emu
        w = int(h * ar)
    return w, h


def add_slides(services_png, jms_png):
    pptx_path = os.path.join(BASE, "PetStore_Architecture_LLD.pptx")
    prs = Presentation(pptx_path) if os.path.exists(pptx_path) else Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW = prs.slide_width
    BLANK = prs.slide_layouts[6]
    NAVY_RGB = RGBColor(0x1F, 0x33, 0x55)
    SUB_RGB = RGBColor(0xC5, 0xD3, 0xE6)

    def header(s, title, subtitle):
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
        band.fill.solid(); band.fill.fore_color.rgb = NAVY_RGB
        band.line.fill.background()
        band.shadow.inherit = False
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.7)).text_frame
        tb.word_wrap = True
        r = tb.paragraphs[0].add_run(); r.text = title
        r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.name = "Calibri"
        sb = s.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.35)).text_frame
        sr = sb.paragraphs[0].add_run(); sr.text = subtitle
        sr.font.size = Pt(13); sr.font.color.rgb = SUB_RGB; sr.font.name = "Calibri"

    def picture_slide(title, subtitle, png):
        s = prs.slides.add_slide(BLANK)
        header(s, title, subtitle)
        avail_w = Inches(12.6); avail_h = Inches(5.95)
        if os.path.exists(png):
            w, h = _fit(png, int(avail_w), int(avail_h))
            x = int((SW - w) / 2); y = Inches(1.28) + int((int(avail_h) - h) / 2)
            s.shapes.add_picture(png, x, y, width=w, height=h)
        else:
            tb = s.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1)).text_frame
            tb.paragraphs[0].add_run().text = f"{os.path.basename(png)} not found — run this script"
        return s

    picture_slide(
        "As-Built Architecture — Service Map",
        "8 Spring Boot services · Service[imported client SDKs] · DB: file H2 (opc/catalog also MongoDB-capable) or none · HTTP=solid, JMS=async",
        services_png)
    picture_slide(
        "As-Built Architecture — JMS Topology",
        "One ActiveMQ Artemis broker · 2 queues + 3 topics + 2 safety nets · producers → destinations → consumers",
        jms_png)

    prs.save(pptx_path)
    print(f"appended 2 slides -> {pptx_path} (now {len(prs.slides)} slides)")


if __name__ == "__main__":
    svc_png = build_services()
    jms_png = build_jms()
    add_slides(svc_png, jms_png)
