#!/usr/bin/env python3
"""
Pet Store — failure-scenario map (as-built resilience behaviour).

Renders one clean Graphviz diagram grouping every failure mode the fleet handles
into six lanes (messaging redelivery, producer outbox, duplicate delivery,
concurrency, business rejection, synchronous HTTP), then appends ONE slide to the
existing deck (docs/PetStore_Architecture_LLD.pptx) non-destructively.

Every value shown is taken from the code:
  - broker.xml: max-delivery-attempts=3, redelivery-delay=1000, multiplier=2.0, max=30000
  - opc.outbox.max-attempts=3 (park), OutboxRelay WARN/ERROR on park
  - inventory processed_event dedup ledger; OPC skip-if-COMPLETED
  - WarehouseOrderEntity @Version optimistic lock -> 409
  - all-or-nothing fulfilment: short stock -> shipped=false, order stays APPROVED
  - admin-office ApiExceptionHandler: 5xx->502, timeout/unreachable->503
  - OPC ApiExceptionHandler: bad input->400, illegal transition/optimistic->409
  - storefront: catalog outage -> "prices unavailable" notice; replayed orderKey -> 400

Outputs: docs/petstore_failure_scenarios.png / .svg  and  +1 slide on the deck.
Run:  python3 docs/petstore_failure_scenarios.py
"""
import os
import graphviz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

BASE = os.path.dirname(os.path.abspath(__file__))

# palette (shared with the deck)
NAVY = "#1F3355"; BLUE = "#2E6DB4"; AMBER = "#C97C2F"; GREEN = "#2F8F46"
RED = "#B23A48"; PURPLE = "#6B4FA0"; GREY = "#555555"; TEAL = "#2A9D8F"


def _card(g, nid, title, body, fill, border, w="2.6"):
    """A labelled step card: coloured title bar + body text.

    Inside Graphviz HTML-like labels a literal ``\\n`` is not a line break — the
    row separator is the ``<BR/>`` tag, so translate before emitting.
    """
    body = body.replace("\\n", "<BR/>")
    g.node(nid, f'<<TABLE BORDER="0" CELLBORDER="1" CELLSPACING="0" CELLPADDING="6">'
                f'<TR><TD BGCOLOR="{border}"><FONT COLOR="white" POINT-SIZE="11"><B>{title}</B></FONT></TD></TR>'
                f'<TR><TD BGCOLOR="{fill}" ALIGN="LEFT"><FONT POINT-SIZE="9.5">{body}</FONT></TD></TR>'
                f'</TABLE>>', shape="plaintext", width=w)


def build():
    g = graphviz.Digraph("failures", format="png")
    g.attr(rankdir="LR", splines="spline", nodesep="0.3", ranksep="0.7", bgcolor="white",
           fontname="Helvetica", fontsize="16", labelloc="t", pad="0.35", compound="true",
           label="Java Pet Store — Failure Scenarios & How the Fleet Recovers\\l"
                 "green = safe outcome / recovery   ·   red = terminal / operator signal   ·   "
                 "every value below is enforced in code\\l")
    g.attr("node", fontname="Helvetica")
    g.attr("edge", fontname="Helvetica", fontsize="9", color=GREY, penwidth="1.3")

    OK = ("#D8F0DD", GREEN)      # safe / recovered
    BAD = ("#F7D9DC", RED)       # terminal / needs attention
    WARN = ("#FBE7D0", AMBER)    # transient / retrying
    INFO = ("#DCE9F7", BLUE)     # trigger / normal

    # ---- Lane 1: JMS consumer throws (transient) --------------------------
    with g.subgraph(name="cluster_msg") as c:
        c.attr(label="1 · Async consumer FAILS to process a message (e.g. DB blip)",
               style="rounded,filled", fillcolor="#F7FAFF", color=BLUE, fontcolor=BLUE, fontsize="12")
        _card(c, "m_trigger", "Consumer throws", "OrderListener / OrderApprovedListener /\\nInvoiceListener raises an exception", *INFO)
        _card(c, "m_retry", "Broker redelivers", "retry 1: ~1s\\nretry 2: ~2s (×2.0, cap 30s)\\nmax-delivery-attempts = 3", *WARN)
        _card(c, "m_ok", "Recovers", "transient fault clears →\\nmessage processed, order proceeds", *OK)
        _card(c, "m_dlq", "3rd failure → DLQ", "poison message routed to DLQ;\\nnotification DlqListener logs ERROR\\n(operator-visible, not lost)", *BAD)
        c.edge("m_trigger", "m_retry")
        c.edge("m_retry", "m_ok", label="succeeds")
        c.edge("m_retry", "m_dlq", label="still failing", color=RED, fontcolor=RED)

    # ---- Lane 2: producer / outbox ----------------------------------------
    with g.subgraph(name="cluster_outbox") as c:
        c.attr(label="2 · OPC cannot PUBLISH an outbound event (broker down)",
               style="rounded,filled", fillcolor="#FFFBF3", color=AMBER, fontcolor=AMBER, fontsize="12")
        _card(c, "o_trigger", "Outbox row unsent", "event written to outbox table in the\\nbusiness tx; OutboxRelay tries to publish", *INFO)
        _card(c, "o_retry", "Relay retries", "@Scheduled re-polls unsent rows;\\nopc.outbox.max-attempts = 3", *WARN)
        _card(c, "o_ok", "Published", "broker back → sent exactly the frozen\\npayload (eventId fixed); marked published", *OK)
        _card(c, "o_park", "Parked (poison)", "3 attempts exhausted → row PARKED,\\nERROR logged, skipped by future polls\\n(no DLQ — it never reached the broker)", *BAD)
        c.edge("o_trigger", "o_retry")
        c.edge("o_retry", "o_ok", label="broker up")
        c.edge("o_retry", "o_park", label="keeps failing", color=RED, fontcolor=RED)

    # ---- Lane 3: duplicate delivery ---------------------------------------
    with g.subgraph(name="cluster_dup") as c:
        c.attr(label="3 · DUPLICATE delivery (JMS is at-least-once)",
               style="rounded,filled", fillcolor="#F3FBF5", color=GREEN, fontcolor=GREEN, fontsize="12")
        _card(c, "d_trigger", "Same event re-delivered", "crash between processing and ACK →\\nbroker re-sends the same eventId", *INFO)
        _card(c, "d_check", "Idempotent consumer", "inventory: processed_event ledger\\n(unique eventId, survives restart)\\nOPC: skip if already COMPLETED", *WARN)
        _card(c, "d_ok", "No double-apply", "replay is a no-op → stock not\\ndecremented twice, order not re-completed", *OK)
        c.edge("d_trigger", "d_check")
        c.edge("d_check", "d_ok", label="seen before → skip")

    # ---- Lane 4: concurrency ----------------------------------------------
    with g.subgraph(name="cluster_conc") as c:
        c.attr(label="4 · CONCURRENT writes (approve + deny race · oversell)",
               style="rounded,filled", fillcolor="#F8F5FC", color=PURPLE, fontcolor=PURPLE, fontsize="12")
        _card(c, "c_trigger", "Two admins act at once", "one approves &amp; one denies the SAME\\nPENDING order simultaneously", *INFO)
        _card(c, "c_lock", "Optimistic + pessimistic lock", "OPC: @Version on the order row\\ninventory: SELECT … FOR UPDATE\\n+ lines sorted by itemId (stable order)", *WARN)
        _card(c, "c_ok", "One wins cleanly", "first commit succeeds; loser gets\\nOptimisticLockingFailure → 409 Conflict.\\nNo oversell, no deadlock", *OK)
        c.edge("c_trigger", "c_lock")
        c.edge("c_lock", "c_ok", label="serialized")

    # ---- Lane 5: business rejection ---------------------------------------
    with g.subgraph(name="cluster_biz") as c:
        c.attr(label="5 · BUSINESS rejections (not errors — designed outcomes)",
               style="rounded,filled", fillcolor="#FDF6F7", color=RED, fontcolor=RED, fontsize="12")
        _card(c, "b_stock", "Short stock", "inventory can't reserve every line", *INFO)
        _card(c, "b_allnone", "All-or-nothing", "NOTHING reserved; InvoiceEvent still\\npublished with shipped=false;\\norder stays APPROVED for later restock", *OK)
        _card(c, "b_key", "Replayed / forged checkout", "missing / reused encrypted orderKey\\n→ 400 invalid_order_key (no duplicate order)", *BAD)
        c.edge("b_stock", "b_allnone")

    # ---- Lane 6: synchronous HTTP -----------------------------------------
    with g.subgraph(name="cluster_http") as c:
        c.attr(label="6 · SYNCHRONOUS HTTP downstream failures",
               style="rounded,filled", fillcolor="#F7FAFF", color=TEAL, fontcolor=TEAL, fontsize="12")
        _card(c, "h_trigger", "Downstream unhealthy", "admin-office calls OPC; storefront\\ncalls catalog/customer", *INFO)
        _card(c, "h_map", "Clean status mapping", "OPC 5xx → 502 Bad Gateway\\nOPC timeout/unreachable → 503\\nbad input → 400 · bad transition → 409", *WARN)
        _card(c, "h_degrade", "Graceful UI degrade", "catalog outage on checkout →\\nrender page + 'prices temporarily\\nunavailable' (no hard 500)", *OK)
        c.edge("h_trigger", "h_map")
        c.edge("h_map", "h_degrade", label="storefront path")

    out = os.path.join(BASE, "petstore_failure_scenarios")
    g.render(out, format="png", cleanup=True)
    g.render(out, format="svg", cleanup=True)
    print("wrote", out + ".png / .svg")
    return out + ".png"


def _fit(img_path, avail_w_emu, avail_h_emu):
    with Image.open(img_path) as im:
        iw, ih = im.size
    ar = iw / ih
    w = avail_w_emu; h = int(w / ar)
    if h > avail_h_emu:
        h = avail_h_emu; w = int(h * ar)
    return w, h


def add_slide(png):
    pptx_path = os.path.join(BASE, "PetStore_Architecture_LLD.pptx")
    prs = Presentation(pptx_path) if os.path.exists(pptx_path) else Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    SW = prs.slide_width
    BLANK = prs.slide_layouts[6]
    NAVY_RGB = RGBColor(0x1F, 0x33, 0x55)
    SUB_RGB = RGBColor(0xC5, 0xD3, 0xE6)

    s = prs.slides.add_slide(BLANK)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY_RGB
    band.line.fill.background(); band.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.7)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = "Failure Scenarios — How the Fleet Recovers"
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r.font.name = "Calibri"
    sb = s.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.35)).text_frame
    sr = sb.paragraphs[0].add_run()
    sr.text = ("6 lanes: async redelivery→DLQ · outbox park · idempotent dedup · "
               "optimistic/pessimistic locks · business rejections · HTTP 502/503 + graceful degrade")
    sr.font.size = Pt(12); sr.font.color.rgb = SUB_RGB; sr.font.name = "Calibri"

    avail_w = Inches(12.7); avail_h = Inches(5.95)
    if os.path.exists(png):
        w, h = _fit(png, int(avail_w), int(avail_h))
        x = int((SW - w) / 2); y = Inches(1.28) + int((int(avail_h) - h) / 2)
        s.shapes.add_picture(png, x, y, width=w, height=h)

    prs.save(pptx_path)
    print(f"appended 1 slide -> {pptx_path} (now {len(prs.slides)} slides)")


if __name__ == "__main__":
    png = build()
    add_slide(png)
